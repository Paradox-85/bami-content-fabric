"""clone_slide: bit-faithful shape copy + image-rel remap."""

from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from shared.pptx.clone import clone_slide, delete_slide_at


def test_clone_preserves_shape_count(template_path):
    prs = Presentation(str(template_path))
    src = prs.slides[1]  # content slide
    n_src_shapes = len(src.shapes)
    n_src_pics = sum(1 for s in src.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

    new, rid_map = clone_slide(prs, src)

    assert len(new.shapes) == n_src_shapes, "shape count must be identical (bit-faithful)"
    n_new_pics = sum(1 for s in new.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert n_new_pics == n_src_pics, "picture count must be identical"
    assert rid_map, "at least one image relationship should be remapped"


def test_clone_background_and_logo_resolve(template_path):
    """The cloned background + logo must still resolve (rels fixed)."""
    prs = Presentation(str(template_path))
    new, _ = clone_slide(prs, prs.slides[1])
    pics = [s for s in new.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # Every picture's blip must reference an existing relationship in the new part.
    from pptx.oxml.ns import qn

    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rel_ids = set(new.part.rels.keys())
    checked = 0
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for pic in pics:
        blips = pic._element.findall(".//" + qn("a:blip"))
        for blip in blips:
            rid = blip.get("{%s}embed") or blip.get("{%s}link" % "")  # noqa
            # Try embed then link explicitly (the r:link attr name):
            rid = blip.get(f"{{{R}}}embed") or blip.get(f"{{{R}}}link")
            if rid is None:
                continue  # effect/decorative blip with no resource reference
            assert rid in rel_ids, f"cloned picture rId {rid} has no matching relationship"
            checked += 1
    assert checked > 0, "no blip references were checked"


def test_delete_slide_removes_from_deck(template_path, tmp_path):
    prs = Presentation(str(template_path))
    n0 = len(prs.slides._sldIdLst)
    clone_slide(prs, prs.slides[1])           # now n0 + 1 slides
    assert len(prs.slides._sldIdLst) == n0 + 1
    for _ in range(n0):
        delete_slide_at(prs, 0)               # prune all originals
    assert len(prs.slides._sldIdLst) == 1
    out = tmp_path / "pruned.pptx"
    prs.save(str(out))
    again = Presentation(str(out))
    assert len(again.slides._sldIdLst) == 1
