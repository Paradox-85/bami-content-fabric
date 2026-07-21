"""Tests for the block-arrow-horizontal native pattern injector.

Verifies:
- Injector is registered and callable
- Resolver can select block-arrow-horizontal via explicit graphical_variant
- Build pipeline works with explicit graphical_variant selection
- Shape naming convention is followed
"""

from __future__ import annotations

import json

from pptx import Presentation

from shared.pptx.pattern_injectors.registry import get_injector, list_injectors
from shared.pptx.pattern_selection import resolve_pattern
from tools.pptx_validate.cli import validate

# ---------------------------------------------------------------------------
# Unit: injector registration
# ---------------------------------------------------------------------------

def test_block_arrow_injector_registered():
    """The block-arrow-horizontal injector is registered and callable."""
    assert "block-arrow-horizontal" in list_injectors()
    fn = get_injector("block-arrow-horizontal")
    assert fn is not None
    assert callable(fn)


# ---------------------------------------------------------------------------
# Unit: resolver enrichment — explicit graphical_variant
# ---------------------------------------------------------------------------

def test_resolver_returns_block_arrow_variant():
    """Content with items + explicit graphical_variant -> block-arrow-horizontal."""
    class FakeTokens:
        def __init__(self, brand="bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    tokens = FakeTokens("bami")
    sel = resolve_pattern(
        content, tokens,
        graphical_variant="block-arrow-horizontal",
    )
    assert sel.family == "numbered-process-steps"
    assert sel.pattern_template_id == "numbered-process-steps/block-arrow-horizontal@1.0.0"
    assert sel.graphical_variant == "block-arrow-horizontal"
    assert sel.renderer_binding is not None
    assert sel.renderer_binding["native"]["injector_id"] == "block-arrow-horizontal"
    assert sel.features is not None
    assert sel.features.get("native_editable") is True


def test_resolver_defaults_to_folded_arrow_without_explicit_variant():
    """Without explicit graphical_variant, resolver returns the first enabled variant (folded-arrow)."""
    class FakeTokens:
        def __init__(self, brand="bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    tokens = FakeTokens("bami")
    sel = resolve_pattern(content, tokens)
    assert sel.family == "numbered-process-steps"
    assert sel.graphical_variant == "folded-arrow-horizontal"


# ---------------------------------------------------------------------------
# Integration: build pipeline with explicit graphical_variant
# ---------------------------------------------------------------------------

def test_build_pipeline_block_arrow_with_explicit_graphical_variant(tmp_path, tmp_out, tokens_path, template_path):
    """A deck with graphical_variant: block-arrow-horizontal resolves and builds."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Block Arrow Build Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Process"},
                "graphical_variant": "block-arrow-horizontal",
                "content": {
                    "items": ["Step A", "Step B", "Step C"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_block_arrow_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: shape naming convention in generated PPTX
# ---------------------------------------------------------------------------

def test_block_arrow_shape_naming_convention(tmp_path, tmp_out, tokens_path, template_path):
    """The generated PPTX contains shapes with the block-arrow naming convention."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Block Arrow Shape Naming Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Named Process"},
                "graphical_variant": "block-arrow-horizontal",
                "content": {
                    "items": ["Alpha", "Beta", "Gamma"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_block_arrow_shape_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    # Open the PPTX and inspect shape names on the content slide (index 1)
    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    shape_names = [shp.name for shp in slide.shapes]
    # Check naming convention pattern
    pattern_shapes = [n for n in shape_names if n.startswith("pattern:numbered-process-steps/block-arrow-horizontal")]
    assert len(pattern_shapes) >= 3, (
        f"Expected at least 3 pattern-named shapes, got {len(pattern_shapes)}: {pattern_shapes}"
    )
    # Verify role naming
    circle_shapes = [n for n in pattern_shapes if n.endswith(":block")]
    number_shapes = [n for n in pattern_shapes if n.endswith(":number")]
    title_shapes = [n for n in pattern_shapes if n.endswith(":title")]
    connector_shapes = [n for n in pattern_shapes if ":connector:" in n]
    assert len(circle_shapes) == 3, f"Expected 3 block shapes, got {len(circle_shapes)}"
    assert len(number_shapes) == 3, f"Expected 3 number shapes, got {len(number_shapes)}"
    assert len(title_shapes) == 3, f"Expected 3 title shapes, got {len(title_shapes)}"
    assert len(connector_shapes) == 2, f"Expected 2 connector shapes, got {len(connector_shapes)}"

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: 5 items (max capacity) builds fine
# ---------------------------------------------------------------------------

def test_block_arrow_5_items(tmp_path, tmp_out, tokens_path, template_path):
    """5 items with block-arrow-horizontal resolve and build successfully."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Block Arrow 5 Steps Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Five Steps"},
                "graphical_variant": "block-arrow-horizontal",
                "content": {
                    "items": ["One", "Two", "Three", "Four", "Five"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_block_arrow_5_steps_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# KVI brand variant tests
# ---------------------------------------------------------------------------


def test_block_arrow_kvi_builds_ok(tmp_path, tmp_out, kvi_template_path, kvi_tokens_path):
    """KVI-branded deck with block-arrow-horizontal resolves and builds."""
    from shared.pptx.build import build_deck
    deck = {
        "title": "KVI Block Arrow Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "KVI Process"},
                "graphical_variant": "block-arrow-horizontal",
                "content": {
                    "items": ["Step A", "Step B", "Step C"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_kvi_block_arrow_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")
    result = build_deck(deck_path, tmp_out, kvi_template_path, kvi_tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    rep = validate(tmp_out, kvi_tokens_path)
    assert rep.ok, f"KVI validation violations: {rep.violations}"
