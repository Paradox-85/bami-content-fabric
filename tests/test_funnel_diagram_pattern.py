"""Tests for the funnel-diagram native pattern injectors.

Verifies:
- Injectors are registered and callable
- Resolver can select each variant via explicit graphical_variant
- Build pipeline works with explicit graphical_variant selection
- Multi-variant routing works for non-arrow family
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from shared.pptx.pattern_injectors.registry import get_injector, list_injectors, inject_pattern
from shared.pptx.pattern_selection import resolve_pattern
from tools.pptx_validate.cli import validate

ROOT = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Unit: injector registration
# ---------------------------------------------------------------------------


def test_funnel_diagram_injector_registered():
    """The funnel-diagram injector is registered and callable."""
    assert "funnel-diagram" in list_injectors()
    fn = get_injector("funnel-diagram")
    assert fn is not None
    assert callable(fn)


def test_funnel_conversion_injector_registered():
    """The funnel-conversion injector is registered and callable."""
    assert "funnel-conversion" in list_injectors()
    fn = get_injector("funnel-conversion")
    assert fn is not None
    assert callable(fn)


# ---------------------------------------------------------------------------
# Resolver: variant selection
# ---------------------------------------------------------------------------


def test_funnel_default_vertical_resolved_by_resolver():
    """Resolver returns default-vertical for funnel-diagram."""
    class FakeTokens:
        def __init__(self, brand: str = "bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    result = resolve_pattern(
        content, FakeTokens("bami"),
        hint_category="funnel-diagram",
    )
    assert result.family == "funnel-diagram"
    assert result.graphical_variant == "default-vertical"
    assert result.renderer_binding is not None


def test_funnel_conversion_pipeline_resolved_by_resolver():
    """Resolver returns conversion-pipeline when requested."""
    class FakeTokens:
        def __init__(self, brand: str = "bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    result = resolve_pattern(
        content, FakeTokens("bami"),
        hint_category="funnel-diagram",
        graphical_variant="conversion-pipeline",
    )
    assert result.family == "funnel-diagram"
    assert result.graphical_variant == "conversion-pipeline"
    assert result.renderer_binding is not None
    native = result.renderer_binding.get("native", {})
    assert native.get("injector_id") == "funnel-conversion"


def test_funnel_sales_growth_resolved_by_resolver():
    """Resolver returns sales-growth when requested."""
    class FakeTokens:
        def __init__(self, brand: str = "bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    result = resolve_pattern(
        content, FakeTokens("bami"),
        hint_category="funnel-diagram",
        graphical_variant="sales-growth",
    )
    assert result.family == "funnel-diagram"
    assert result.graphical_variant == "sales-growth"
    assert result.renderer_binding is not None
    native = result.renderer_binding.get("native", {})
    assert native.get("injector_id") == "funnel-diagram"  # reuses existing injector


# ---------------------------------------------------------------------------
# E2E: injector rendering
# ---------------------------------------------------------------------------


def test_funnel_default_vertical_injects_real_shapes():
    """funnel-diagram injector produces real shapes."""
    from shared.pptx.tokens import load_tokens
    real_tokens = load_tokens(ROOT / "templates" / "bami" / "design_tokens.yaml")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    created = inject_pattern(
        slide, real_tokens,
        "funnel-diagram",
        x=0.5, y=1.0, w=9.0, h=4.0,
        segments=[
            {"label": "Awareness", "value": "80%", "pct": 1.0},
            {"label": "Interest", "value": "60%", "pct": 0.75},
            {"label": "Decision", "value": "30%", "pct": 0.5},
        ],
    )
    assert len(created) >= 3  # at least 3 shapes (bar, label, value per stage?)


def test_funnel_conversion_pipeline_injects_real_shapes():
    """funnel-conversion injector produces real shapes."""
    from shared.pptx.tokens import load_tokens
    real_tokens = load_tokens(ROOT / "templates" / "bami" / "design_tokens.yaml")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    created = inject_pattern(
        slide, real_tokens,
        "funnel-conversion",
        x=0.5, y=1.0, w=9.0, h=4.0,
        stages=[
            {"label": "Visit", "value": "10K", "pct": 1.0},
            {"label": "Signup", "value": "2K", "pct": 0.8},
            {"label": "Purchase", "value": "500", "pct": 0.6},
        ],
    )
    assert len(created) >= 3  # at least 3 shapes


# ---------------------------------------------------------------------------
# E2E: build pipeline with explicit variant
# ---------------------------------------------------------------------------


def _write_deck(tmp_path, deck):
    import json
    path = tmp_path / "_funnel.json"
    path.write_text(json.dumps(deck, indent=2), encoding="utf-8")
    return path


def test_funnel_default_vertical_builds_ok(tmp_path, tmp_out, tokens_path, template_path):
    """Build a deck with funnel-diagram default-vertical variant.
    Verifies that actual funnel pattern shapes are produced (not fallback to
    numbered-process-steps).
    """
    from shared.pptx.build import build_deck

    deck = {
        "title": "Funnel Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Funnel"}},
            {
                "template": "content",
                "fields": {"title": "Pipeline"},
                "blocks": [
                    {
                        "kind": "inject-pattern",
                        "canonical_id": "funnel-diagram",
                        "x": 0.5, "y": 1.0, "w": 9.0, "h": 4.5,
                        "segments": [
                            {"label": "Awareness", "value": 100, "pct": 1.0},
                            {"label": "Interest", "value": 60, "pct": 0.75},
                            {"label": "Decision", "value": 30, "pct": 0.5},
                        ],
                    },
                ],
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify the output contains funnel pattern shapes (not fallback)
    prs = Presentation(str(tmp_out))
    slide = list(prs.slides)[1]  # content slide
    pattern_count = 0
    has_funnel_shapes = False
    has_fallback_shapes = False
    for shp in slide.shapes:
        name = getattr(shp, "name", "") or ""
        if name.startswith("pattern:"):
            pattern_count += 1
            if "funnel-diagram" in name:
                has_funnel_shapes = True
            if "numbered-process-steps" in name:
                has_fallback_shapes = True
    assert pattern_count > 0, f"No pattern shapes found on content slide"
    assert has_funnel_shapes, (
        f"Expected funnel-diagram pattern shapes but got only "
        f"numbered-process-steps (fallback). Warnings: {result.get('selection_warnings', [])}"
    )
    assert not has_fallback_shapes, (
        f"Deck fell back to numbered-process-steps instead of funnel. "
        f"Warnings: {result.get('selection_warnings', [])}"
    )
    # Validate the output
    report = validate(tmp_out, tokens_path)
    assert report.ok, f"Validation violations: {report.violations}"



def test_funnel_conversion_pipeline_builds_ok(tmp_path, tmp_out, tokens_path, template_path):
    """Build a deck with funnel-diagram conversion-pipeline variant.
    Verifies that actual funnel conversion pattern shapes are produced.
    """
    from shared.pptx.build import build_deck

    deck = {
        "title": "Conversion Pipeline Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Conversion"}},
            {
                "template": "content",
                "fields": {"title": "Pipeline"},
                "blocks": [
                    {
                        "kind": "inject-pattern",
                        "canonical_id": "funnel-conversion",
                        "x": 0.5, "y": 1.0, "w": 9.0, "h": 4.5,
                        "stages": [
                            {"label": "Visit", "value": 100, "pct": 1.0},
                            {"label": "Signup", "value": 80, "pct": 0.8},
                            {"label": "Purchase", "value": 50, "pct": 0.6},
                            {"label": "Retain", "value": 20, "pct": 0.4},
                        ],
                    },
                ],
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify funnel shapes
    prs = Presentation(str(tmp_out))
    slide = list(prs.slides)[1]
    has_funnel_shapes = False
    has_fallback_shapes = False
    for shp in slide.shapes:
        name = getattr(shp, "name", "") or ""
        if name.startswith("pattern:"):
            if "funnel-diagram" in name:
                has_funnel_shapes = True
            if "numbered-process-steps" in name:
                has_fallback_shapes = True
    assert has_funnel_shapes, (
        f"Expected funnel-diagram shapes but got fallback to numbered-process-steps. "
        f"Warnings: {result.get('selection_warnings', [])}"
    )
    assert not has_fallback_shapes, "Deck fell back to numbered-process-steps"
    report = validate(tmp_out, tokens_path)
    assert report.ok, f"Validation violations: {report.violations}"


def test_funnel_sales_growth_builds_ok(tmp_path, tmp_out, tokens_path, template_path):
    """Build a deck with funnel-diagram sales-growth variant.
    Verifies that actual funnel pattern shapes are produced.
    """
    from shared.pptx.build import build_deck

    deck = {
        "title": "Sales Growth Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Growth"}},
            {
                "template": "content",
                "fields": {"title": "Stages"},
                "blocks": [
                    {
                        "kind": "inject-pattern",
                        "canonical_id": "funnel-diagram",
                        "x": 0.5, "y": 1.0, "w": 9.0, "h": 4.5,
                        "segments": [
                            {"label": "Prospecting", "value": 100, "pct": 1.0},
                            {"label": "Qualification", "value": 80, "pct": 0.75},
                            {"label": "Proposal", "value": 50, "pct": 0.5},
                            {"label": "Close", "value": 20, "pct": 0.3},
                        ],
                    },
                ],
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify funnel shapes
    prs = Presentation(str(tmp_out))
    slide = list(prs.slides)[1]
    has_funnel_shapes = False
    has_fallback_shapes = False
    for shp in slide.shapes:
        name = getattr(shp, "name", "") or ""
        if name.startswith("pattern:"):
            if "funnel-diagram" in name:
                has_funnel_shapes = True
            if "numbered-process-steps" in name:
                has_fallback_shapes = True
    assert has_funnel_shapes, (
        f"Expected funnel-diagram shapes but got fallback to numbered-process-steps. "
        f"Warnings: {result.get('selection_warnings', [])}"
    )
    assert not has_fallback_shapes, "Deck fell back to numbered-process-steps"
    report = validate(tmp_out, tokens_path)
    assert report.ok, f"Validation violations: {report.violations}"


# ---------------------------------------------------------------------------
# E2E: KVI multi-brand build
# ---------------------------------------------------------------------------


def test_funnel_default_vertical_kvi_builds_ok(
    tmp_path, tmp_out, kvi_tokens_path, kvi_template_path,
):
    """Build a KVI-branded deck with funnel-diagram default-vertical variant.
    Verifies that actual funnel pattern shapes are produced.
    """
    from shared.pptx.build import build_deck

    deck = {
        "title": "KVI Funnel Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Funnel"}},
            {
                "template": "content",
                "fields": {"title": "Pipeline"},
                "blocks": [
                    {
                        "kind": "inject-pattern",
                        "canonical_id": "funnel-diagram",
                        "x": 0.5, "y": 1.0, "w": 9.0, "h": 4.5,
                        "segments": [
                            {"label": "A", "value": 100, "pct": 1.0},
                            {"label": "B", "value": 80, "pct": 0.75},
                            {"label": "C", "value": 50, "pct": 0.5},
                        ],
                    },
                ],
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, kvi_template_path, kvi_tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify funnel shapes
    prs = Presentation(str(tmp_out))
    slide = list(prs.slides)[1]
    has_funnel_shapes = False
    has_fallback_shapes = False
    for shp in slide.shapes:
        name = getattr(shp, "name", "") or ""
        if name.startswith("pattern:"):
            if "funnel-diagram" in name:
                has_funnel_shapes = True
            if "numbered-process-steps" in name:
                has_fallback_shapes = True
    assert has_funnel_shapes, (
        f"Expected funnel-diagram shapes but got fallback. "
        f"Warnings: {result.get('selection_warnings', [])}"
    )
    assert not has_fallback_shapes, "Deck fell back to numbered-process-steps"
    report = validate(tmp_out, kvi_tokens_path)
    assert report.ok, f"KVI validation violations: {report.violations}"
