"""Tests for the circle-steps (numbered circle-step diagram) native PPTX injector."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from shared.pptx.tokens import load_tokens

ROOT = Path(__file__).resolve().parent.parent


def _load_tokens(tokens_path: Path):
    """Helper to load tokens from a path."""
    return load_tokens(tokens_path)


class TestCircleStepsInjector:
    def test_circle_steps_injector_creates_shapes(self, tmp_path, tokens_path):
        """Circle-steps injector should produce shapes for each node."""
        tokens = _load_tokens(tokens_path)

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("circle-steps")
        assert injector is not None, "circle-steps injector not registered"

        nodes = [
            {"label": "Research"},
            {"label": "Design"},
            {"label": "Build"},
            {"label": "Test"},
        ]

        shapes = injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0, nodes=nodes)
        # Should create: 4 connectors + 4 circles + 4 numbers + 4 labels = 16+
        assert len(shapes) >= 8, f"Expected >=8 shapes, got {len(shapes)}"

    def test_circle_steps_injector_raises_without_nodes(self, tmp_path, tokens_path):
        """Circle-steps injector should raise ValueError when nodes is missing."""
        tokens = _load_tokens(tokens_path)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("circle-steps")
        assert injector is not None

        with pytest.raises(ValueError, match="circle-steps: 'nodes' parameter is required"):
            injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0)

    def test_circle_steps_handles_3_nodes(self, tmp_path, tokens_path):
        """Circle-steps should work with 3 nodes."""
        tokens = _load_tokens(tokens_path)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("circle-steps")
        assert injector is not None

        nodes = [
            {"label": "Step 1"},
            {"label": "Step 2"},
            {"label": "Step 3"},
        ]

        shapes = injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0, nodes=nodes)
        assert len(shapes) >= 6, f"Expected >=6 shapes, got {len(shapes)}"

    def test_circle_steps_handles_6_nodes(self, tmp_path, tokens_path):
        """Circle-steps should work with 6 nodes (max items)."""
        tokens = _load_tokens(tokens_path)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("circle-steps")
        assert injector is not None

        nodes = [
            {"label": f"Step {i+1}"} for i in range(6)
        ]

        shapes = injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0, nodes=nodes)
        assert len(shapes) >= 12, f"Expected >=12 shapes, got {len(shapes)}"

    def test_circle_steps_variant_registered_in_registry(self):
        """The circular-process-loop/circle-steps variant should be registered in the registry."""
        import yaml

        registry_path = ROOT / "schemas" / "pattern-registry.yaml"
        with registry_path.open(encoding="utf-8") as f:
            registry = yaml.safe_load(f)

        for entry in registry.get("entries", []):
            if entry.get("family") == "circular-process-loop":
                variants = [v.get("graphical_variant") for v in entry.get("graphical_variants", [])]
                assert "circle-steps" in variants, (
                    "circle-steps variant not found in circular-process-loop family"
                )
                return
        pytest.fail("circular-process-loop family not found in registry")