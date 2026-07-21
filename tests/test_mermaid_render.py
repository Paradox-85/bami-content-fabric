"""Tests for shared.pptx.mermaid_render (T6)."""

from __future__ import annotations

import subprocess
from pathlib import Path

import pytest
from PIL import Image

from shared.pptx.mermaid_render import (
    MermaidRenderError,
    mmdc_available,
    render_mermaid_png,
)
from tests.conftest import ROOT

# ---------------------------------------------------------------------------
# Module-level skip detection
# ---------------------------------------------------------------------------

_HAVE_MMDC = mmdc_available()
_SAMPLE = "flowchart LR\n  A[One] --> B[Two]\n  B --> C[Three]"


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

@pytest.mark.skipif(not _HAVE_MMDC, reason="mmdc not installed")
class TestWithMmdc:
    """Integration tests that require the actual mmdc binary."""

    def test_render_mermaid_png_produces_valid_png(self):
        """Render a simple flowchart and verify the output is a real PNG."""
        path = render_mermaid_png(_SAMPLE)
        assert path.exists(), f"PNG not found at {path}"
        with Image.open(path) as im:
            assert im.format == "PNG"
            assert im.size[0] > 0
            assert im.size[1] > 0

    def test_cache_hit_skips_rerender(self):
        """Calling render_mermaid_png twice with the same definition reuses
        the cached file without invoking mmdc again."""
        import shared.pptx.mermaid_render as mr

        unique_def = f"flowchart LR\n  A[UniqueTest] --> B[{id(self)}]\n"

        # First call forces a real render.
        path1 = render_mermaid_png(unique_def)

        # Patch subprocess.run to count invocations.
        original_run = mr.subprocess.run
        call_count = 0

        def counting_run(*args, **kwargs):
            nonlocal call_count
            call_count += 1
            return original_run(*args, **kwargs)

        mr.subprocess.run = counting_run
        try:
            path2 = render_mermaid_png(unique_def)
        finally:
            mr.subprocess.run = original_run

        assert path1 == path2, "Cache should return the same path"
        assert call_count == 0, (
            f"subprocess.run was called {call_count} time(s); "
            "expected 0 on cache hit"
        )


class TestWithoutMmdc:
    """Tests that monkeypatch the mmdc binary / subprocess to avoid a real
    toolchain dependency. These pass even without npm install."""

    def test_mmdc_missing_raises_loud(self, monkeypatch):
        """When mmdc is absent, MermaidRenderError is raised."""
        monkeypatch.setattr(
            "shared.pptx.mermaid_render._mmdc_argv",
            lambda: None,
        )
        with pytest.raises(MermaidRenderError) as excinfo:
            render_mermaid_png("A-->B")
        assert "npm install" in str(excinfo.value)

    def test_render_error_raises_loud(self, monkeypatch):
        """When mmdc returns non-zero, MermaidRenderError is raised."""
        def _fake_argv():
            return ["mmdc_fake"]

        monkeypatch.setattr(
            "shared.pptx.mermaid_render._mmdc_argv", _fake_argv
        )

        def _fail_run(*args, **kwargs):
            return subprocess.CompletedProcess(
                args=[], returncode=1, stdout="", stderr="boom"
            )

        monkeypatch.setattr(
            "shared.pptx.mermaid_render.subprocess.run", _fail_run
        )
        with pytest.raises(MermaidRenderError) as excinfo:
            render_mermaid_png("A-->B")
        assert "boom" in str(excinfo.value)


@pytest.mark.skipif(not _HAVE_MMDC, reason="mmdc not installed")
class TestIntegration:
    """End-to-end: build a deck with a Mermaid image block, then validate."""

    @pytest.mark.xfail(strict=False, reason="deferred E6: 'image' block kind not in schema enum/BUILDERS; see docs/runbooks/library-runtime-error-log.md")
    def test_mermaid_image_block_builds_and_validates(
        self, tmp_out, tokens_path,
    ):
        """Build a 3-slide deck whose content block is a Mermaid image and
        validate it."""
        from shared.pptx.build import build_deck
        from shared.pptx.schema import validate_deck

        deck = {
            "schema_version": 2,
            "title": "Mermaid test deck",
            "slides": [
                {"template": "cover", "fields": {}},
                {
                    "template": "content",
                    "fields": {"title": "Test"},
                    "blocks": [
                        {
                            "kind": "image",
                            "x": 1.5,
                            "y": 2.4,
                            "w": 16.8,
                            "h": 7.2,
                            "fit": "contain",
                            "src": {"mermaid": _SAMPLE},
                        }
                    ],
                },
                {"template": "closing", "fields": {}},
            ],
        }

        # Write a temp deck.json so load_deck can work with its path-based
        # token resolution.
        deck_dir = Path(tmp_out).parent
        deck_path = deck_dir / "_test_mermaid_deck.json"
        import json
        deck_path.write_text(json.dumps(deck), encoding="utf-8")

        # Validate schema first
        validate_deck(deck)

        template_path = ROOT / "templates" / "template.pptx"

        # Build
        build_deck(str(deck_path), str(tmp_out), str(template_path), str(tokens_path))

        # Validate
        from tools.pptx_validate.cli import validate

        rep = validate(str(tmp_out), str(tokens_path))
        assert rep.ok, (
            f"Validator violations: {rep.violations}"
        )

        # Verify the Mermaid image was embedded as a body-zone picture.
        # Chrome contributes a full-bleed background at (0,0) + a small logo;
        # the Mermaid diagram is the only picture placed inside the declared
        # body zone (x=1.5", y=2.4", w=16.8", h=7.2").
        from pptx import Presentation

        prs = Presentation(str(tmp_out))
        content_slide = prs.slides[1]
        pic_shapes = [s for s in content_slide.shapes if s.shape_type == 13]  # 13 = Picture
        EMU_IN = 914400
        body_pics = [
            p for p in pic_shapes
            if p.left is not None and p.top is not None
            and p.left > 1.2 * EMU_IN           # not the (0,0) full-bleed background
            and p.top > 2.0 * EMU_IN            # below the title-bar / logo row
            and (p.width or 0) > 2.0 * EMU_IN   # a real diagram, not a small icon
        ]
        assert body_pics, (
            "No Mermaid picture embedded in the body zone; content-slide picture "
            f"count = {len(pic_shapes)} (chrome-only baseline is ~2)"
        )
        mermaid = body_pics[0]
        # Declared zone starts at x=1.5in, y=2.4in; contain-fit keeps the picture inside it.
        assert abs(mermaid.left - 1.5 * EMU_IN) < 0.6 * EMU_IN, f"left={mermaid.left}"
        assert mermaid.top >= int(2.3 * EMU_IN), f"top={mermaid.top}"
