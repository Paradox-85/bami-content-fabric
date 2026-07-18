"""Validator: passes on a clean deck and flags seeded brand violations."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

from shared.pptx.build import build_deck
from tools.pptx_validate.cli import validate


def _build(sample_deck, tokens_path, template_path, tmp_out):
    build_deck(sample_deck, tmp_out, template_path, tokens_path)
    return Presentation(str(tmp_out))


def _save(prs, path):
    prs.save(str(path))
    return path


def test_clean_deck_passes(sample_deck, tokens_path, template_path, tmp_out):
    build_deck(sample_deck, tmp_out, template_path, tokens_path)
    rep = validate(tmp_out, tokens_path)
    assert rep.ok


def test_flags_non_montserrat_font(sample_deck, tokens_path, template_path, tmp_out):
    prs = _build(sample_deck, tokens_path, template_path, tmp_out)
    # Mutate one body run on a content slide to a non-brand font.
    sl = prs.slides[1]
    for shp in sl.shapes:
        if shp.has_text_frame and shp.text_frame.paragraphs:
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name == "Montserrat":
                        r.font.name = "Arial"
                        break
            else:
                continue
            break
    _save(prs, tmp_out)
    rep = validate(tmp_out, tokens_path)
    assert any("not Montserrat" in v for v in rep.violations), rep.violations


def test_flags_off_brand_color(sample_deck, tokens_path, template_path, tmp_out):
    prs = _build(sample_deck, tokens_path, template_path, tmp_out)
    sl = prs.slides[1]
    for shp in sl.shapes:
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        r.font.color.rgb = RGBColor(0xFF, 0x00, 0xFF)  # magenta
                        _save(prs, tmp_out)
                        rep = validate(tmp_out, tokens_path)
                        assert any("outside the brand palette" in v for v in rep.violations), rep.violations
                        return
    raise AssertionError("no mutable run found")


def test_flags_out_of_bounds(sample_deck, tokens_path, template_path, tmp_out):
    prs = _build(sample_deck, tokens_path, template_path, tmp_out)
    sl = prs.slides[1]
    # Add a textbox far off the right edge of the canvas.
    sl.shapes.add_textbox(Inches(30), Inches(2), Inches(2), Inches(1))
    _save(prs, tmp_out)
    rep = validate(tmp_out, tokens_path)
    assert any("out of canvas bounds" in v for v in rep.violations), rep.violations


def test_flags_missing_logo(sample_deck, tokens_path, template_path, tmp_out):
    prs = _build(sample_deck, tokens_path, template_path, tmp_out)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    sl = prs.slides[1]
    # Remove the BAMI logo picture (top-right).
    for shp in list(sl.shapes):
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and shp.left and shp.left > 15 * 914400:
            shp._element.getparent().remove(shp._element)
    _save(prs, tmp_out)
    rep = validate(tmp_out, tokens_path)
    assert any("brand logo not at the token EMU position" in v for v in rep.violations), rep.violations
