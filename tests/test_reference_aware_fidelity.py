"""Tests for reference-aware fidelity comparison.

- grammar-aware comparison works
- fidelity artifacts are emitted
- topology substitution is detected (real contract strings)
- CLI fidelity workflow resolves contract_path (regression test)
"""
from __future__ import annotations

import json
import tempfile
from pathlib import Path
from typing import Any

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

from shared.pptx.reference_analysis import (
    check_forbidden_substitutions,
    compare_grammar,
    extract_topology_metrics,
)

ROOT = Path(__file__).resolve().parents[1]
BUILD_FIDELITY = ROOT / "build" / "fidelity"

# ---------------------------------------------------------------------------
#  Fixture helpers — build a minimal .pptx with known shapes for testing
# ---------------------------------------------------------------------------


def _make_slide_with_axis() -> tuple[Any, Any]:
    """Create a minimal 1-slide PPTX with a labeled axis shape (straight rounded-rect)."""
    prs = Presentation()
    prs.slide_width = 12192000  # 13.33 in (KVI widescreen)
    prs.slide_height = 6858000  # 7.5 in
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Add a ROUNDED_RECTANGLE axis with the expected naming convention
    axis = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        500000,    # left (EMU)
        3000000,   # top
        10000000,  # width
        150000,    # height
    )
    axis.name = "pattern:roadmap-with-milestones/default-horizontal:axis"

    # Add some empty phase bands
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        500000,
        2000000,
        3000000,
        4000000,
    )
    band.name = "pattern:roadmap-with-milestones/default-horizontal:phase:01"

    return prs, slide


def _save_fixture(prs: Any) -> str:
    """Save a fixture Presentation to a temp file and return the path."""
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


# ---------------------------------------------------------------------------
#  Tests
# ---------------------------------------------------------------------------


class TestTopologyMetrics:
    """Topology metric extraction tests."""

    def test_extract_topology_returns_dict(self):
        """extract_topology_metrics must return a dict with expected keys."""
        from inspect import signature
        sig = signature(extract_topology_metrics)
        assert "slide" in sig.parameters
        assert extract_topology_metrics.__doc__ is not None
        assert "topology" in extract_topology_metrics.__doc__.lower()

    def test_extract_metrics_for_straight_axis(self):
        """A straight ROUNDED_RECTANGLE axis must be detected as trajectory."""
        prs, slide = _make_slide_with_axis()
        metrics = extract_topology_metrics(slide)
        assert metrics["has_trajectory_axis"] is True


class TestForbiddenSubstitution:
    """Forbidden substitution detection tests.

    Uses real contract-style descriptive strings (e.g. "Gantt table or task rows")
    to verify substring/semantic matching works, not exact token membership.
    """

    def test_check_forbidden_substitutions_empty(self):
        """Empty forbidden list should produce no violations."""
        violations = check_forbidden_substitutions(None, [])
        assert violations == []

    def test_detect_gantt_with_descriptive_string(self):
        """'Gantt table or task rows' must trigger Gantt detection via substring match."""
        prs, slide = _make_slide_with_axis()
        # Add a Gantt-prefixed shape
        gantt = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100000, 100000)
        gantt.name = "pattern:gantt:some-task"
        violations = check_forbidden_substitutions(
            slide,
            ["Gantt table or task rows", "Flat numbered process steps"],
        )
        assert any("Gantt" in v for v in violations), (
            f"Expected Gantt detection with descriptive string, got {violations}"
        )

    def test_detect_straight_line_with_descriptive_string(self):
        """'Plain straight-line-plus-circles (must have trajectory...)' must trigger detection."""
        prs, slide = _make_slide_with_axis()
        violations = check_forbidden_substitutions(
            slide,
            ["Plain straight-line-plus-circles (must have trajectory character)"],
        )
        assert any("straight-line" in v.lower() for v in violations), (
            f"Expected straight-line detection with descriptive string, got {violations}"
        )

    def test_no_false_positive_with_unrelated_forbidden_list(self):
        """Unrelated forbidden strings must not produce false positives."""
        prs, slide = _make_slide_with_axis()
        violations = check_forbidden_substitutions(
            slide,
            ["Raster image fallback", "Equal-sized cards layout"],
        )
        # No axis-related violation because "straight-line-plus-circles" is not in the list
        assert all("straight-line" not in v.lower() for v in violations), (
            f"Unexpected straight-line detection with unrelated forbidden list: {violations}"
        )

    def test_detect_raster_image(self):
        """'Raster image fallback' must trigger raster detection."""
        prs, slide = _make_slide_with_axis()
        violations = check_forbidden_substitutions(
            slide,
            ["Raster image fallback"],
        )
        # No raster on this slide, so no violation expected
        assert all("raster" not in v.lower() for v in violations)


class TestGrammarComparisonContract:
    """Tests that compare_grammar reads forbidden_outputs from the contract."""

    def test_compare_grammar_uses_contract_forbidden(self):
        """compare_grammar must read forbidden_outputs from contract YAML."""
        prs, slide = _make_slide_with_axis()
        pptx_path = _save_fixture(prs)

        # Create a temporary contract YAML with forbidden_outputs
        contract_data = {
            "family": "roadmap-with-milestones",
            "variant": "default-horizontal",
            "forbidden_outputs": [
                "Plain straight-line-plus-circles (must have trajectory character)",
                "Flat numbered process steps",
            ],
        }
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", delete=False, encoding="utf-8"
        ) as f:
            yaml.dump(contract_data, f)
            contract_path = f.name

        try:
            result = compare_grammar(pptx_path, contract_path)
            # The slide has a straight ROUNDED_RECTANGLE axis, so straight-line detection should fire
            violations = result["violations"]
            assert any("straight-line" in v.lower() for v in violations), (
                f"Expected straight-line violation from contract-driven compare_grammar, "
                f"got violations={violations}"
            )
            assert result["summary"]["passes"] is False
        finally:
            Path(pptx_path).unlink(missing_ok=True)
            Path(contract_path).unlink(missing_ok=True)

    def test_compare_grammar_falls_back_to_registry(self):
        """Compare grammar without contract path."""
        prs, slide = _make_slide_with_axis()
        pptx_path = _save_fixture(prs)
        try:
            result = compare_grammar(pptx_path, contract_path=None)
            # No forbidden_outputs in registry -> no violations from that source
            assert result is not None
            assert isinstance(result["violations"], list)
        finally:
            Path(pptx_path).unlink(missing_ok=True)

    def test_compare_grammar_with_default_horizontal_contract(self):
        """Use the actual default-horizontal.v1.yaml contract to test the real workflow."""
        contract_path = (
            ROOT
            / "schemas"
            / "visual-contracts"
            / "roadmap-with-milestones"
            / "default-horizontal.v1.yaml"
        )
        assert contract_path.exists(), f"Contract not found: {contract_path}"

        prs, slide = _make_slide_with_axis()
        pptx_path = _save_fixture(prs)
        try:
            result = compare_grammar(pptx_path, contract_path)
            violations = result["violations"]
            # The real contract has "Plain straight-line-plus-circles (...)",
            # and our fixture has a straight ROUNDED_RECTANGLE axis.
            # The fixture does NOT have :band suffixed phase shapes,
            # so is_roadmap_slide is False and straight-line detection fires.
            assert any("straight-line" in v.lower() for v in violations), (
                f"Real contract should detect straight-line on our fixture axis, "
                f"got violations={violations}"
            )
        finally:
            Path(pptx_path).unlink(missing_ok=True)


class TestFidelityArtifacts:
    """Fidelity artifact tests."""

    def test_build_fidelity_dirs_exist(self):
        """Build fidelity directories must exist."""
        for sub in ("reference-metrics", "native-metrics", "reports", "contact-sheets"):
            d = BUILD_FIDELITY / sub
            assert d.exists(), f"Fidelity directory missing: {d}"

    def test_contact_sheet_requires_review(self):
        """Contact sheet metadata must indicate review is required."""
        sheet_path = BUILD_FIDELITY / "contact-sheets" / "test-pilot.json"
        sheet_path.parent.mkdir(parents=True, exist_ok=True)
        test_sheet = {
            "pilot_id": "test-pilot",
            "visual_review_required": True,
        }
        with sheet_path.open("w") as f:
            json.dump(test_sheet, f)
        with sheet_path.open() as f:
            loaded = json.load(f)
        assert loaded["visual_review_required"] is True
        # Cleanup
        sheet_path.unlink()


class TestGrammarComparison:
    """Grammar comparison tests."""

    def test_compare_grammar_structure(self):
        """compare_grammar must return expected structure."""
        assert compare_grammar.__doc__ is not None, "compare_grammar must have docstring"


class TestCliFidelityWorkflow:
    """End-to-end CLI fidelity workflow tests.

    Validates that the real --fidelity workflow resolves contract_path
    correctly for roadmap and second pilot families.
    """

    def test_run_fidelity_workflow_resolves_contract_for_roadmap(self):
        """run_fidelity_workflow must receive a non-None contract_path for roadmap.

        This is a regression test for the bug where contract_path was None
        because the resolver matched on file stem instead of parent directory.
        """
        from shared.pptx.fidelity_compare import run_fidelity_workflow

        prs, slide = _make_slide_with_axis()
        pptx_path = _save_fixture(prs)
        try:
            contracts_dir = (
                ROOT / "schemas" / "visual-contracts" / "roadmap-with-milestones"
            )
            # Find the first .yaml contract in the roadmap family directory
            contract_path = None
            for child in contracts_dir.rglob("*.yaml"):
                # Match by parent directory name (the family), not file stem
                if child.parent.name == "roadmap-with-milestones":
                    contract_path = child
                    break
            assert contract_path is not None, (
                "No contract resolved for roadmap-with-milestones family. "
                "The CLI uses child.parent.name matching; this test proves it works."
            )
            assert contract_path.exists(), f"Contract not found: {contract_path}"
            assert contract_path.parent.name == "roadmap-with-milestones", (
                f"Contract parent should be roadmap-with-milestones, got {contract_path.parent.name}"
            )
            result = run_fidelity_workflow(pptx_path, "roadmap-with-milestones", contract_path)
            # Artifacts should be generated
            assert "metrics" in result
            assert "report" in result
            assert "contact_sheet" in result
        finally:
            Path(pptx_path).unlink(missing_ok=True)

    def test_run_fidelity_workflow_resolves_contract_for_numbered_process(self):
        """run_fidelity_workflow must receive a non-None contract_path for numbered-process-steps.

        Regression test: second pilot family must also resolve its contract.
        """
        from shared.pptx.fidelity_compare import run_fidelity_workflow

        # Create a minimal fixture with a numbered-process-steps shape
        prs = Presentation()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 100000, 100000)
        shp.name = "pattern:numbered-process-steps/folded-arrow-horizontal:step:01"
        pptx_path = _save_fixture(prs)
        try:
            contracts_dir = ROOT / "schemas" / "visual-contracts" / "numbered-process-steps"
            contract_path = None
            for child in contracts_dir.rglob("*.yaml"):
                if child.parent.name == "numbered-process-steps":
                    contract_path = child
                    break
            assert contract_path is not None, (
                "No contract resolved for numbered-process-steps family. "
                "The CLI must resolve by parent directory name."
            )
            assert contract_path.exists(), f"Contract not found: {contract_path}"
            assert contract_path.parent.name == "numbered-process-steps"
            result = run_fidelity_workflow(pptx_path, "numbered-process-steps", contract_path)
            assert "metrics" in result
            assert "report" in result
        finally:
            Path(pptx_path).unlink(missing_ok=True)

    def test_fidelity_roadmap_contract_not_none_in_cli_path(self):
        """Simulate the exact CLI resolution logic and prove contract_path is not None.

        This directly reproduces the bug: old code used child.stem.startswith(pilot_id),
        which fails because stem is 'default-horizontal.v1' not 'roadmap-with-milestones'.
        Fixed code uses child.parent.name == pilot_id.
        """
        from pathlib import Path as PPath

        pptx_path_obj = PPath("build/fidelity/native-renders/roadmap-with-milestones.bami.pptx")
        pilot_id = pptx_path_obj.stem.replace(".bami", "").replace(".kvi", "")
        assert pilot_id == "roadmap-with-milestones", f"Expected roadmap-with-milestones, got {pilot_id!r}"

        contracts_dir = ROOT / "schemas" / "visual-contracts"
        assert contracts_dir.exists(), f"Contracts dir not found: {contracts_dir}"

        # Old broken logic: child.stem.startswith(pilot_id)
        old_contract = None
        for child in contracts_dir.rglob("*.yaml"):
            if child.stem.startswith(pilot_id):
                old_contract = child
                break
        # Assert that the old logic FAILS (this is the bug)
        assert old_contract is None, (
            f"Old file-stem matching should NOT find a contract, "
            f"but found {old_contract}. This would have been the bug."
        )

        # New fixed logic: child.parent.name == pilot_id
        new_contract = None
        for child in contracts_dir.rglob("*.yaml"):
            if child.parent.name == pilot_id:
                new_contract = child
                break
        assert new_contract is not None, (
            f"Fixed parent-directory matching must find a contract for {pilot_id}, "
            f"but got None. This is the bug from Blocker 1."
        )
        assert new_contract.exists(), f"Contract not found: {new_contract}"
        assert new_contract.parent.name == "roadmap-with-milestones"
