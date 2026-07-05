"""chrome.apply_slots: text replaced, run formatting preserved."""

from __future__ import annotations

from pptx import Presentation

from shared.pptx.chrome import apply_slots, set_slot_text, shape_by_name
from shared.pptx.clone import clone_slide


def test_set_slot_preserves_run_format(template_path):
    prs = Presentation(str(template_path))
    new, _ = clone_slide(prs, prs.slides[1])  # content slide
    title_shape = shape_by_name(new, "Text 1")
    assert title_shape is not None
    # Before: original title text
    assert title_shape.text_frame.text.strip() == "Context & proposal"
    # Replace
    assert set_slot_text(title_shape, "Quarterly review")
    # After: new text, formatting preserved
    r = title_shape.text_frame.paragraphs[0].runs[0]
    assert r.text == "Quarterly review"
    assert r.font.name == "Montserrat"
    assert r.font.size.pt == 24.0
    assert r.font.bold is True
    assert str(r.font.color.rgb) == "FFFFFF"


def test_apply_slots_fills_content_title(template_path):
    prs = Presentation(str(template_path))
    new, _ = clone_slide(prs, prs.slides[1])
    slots = {"title": "Text 1"}
    apply_slots(new, slots, {"title": "Use cases by department"})
    assert shape_by_name(new, "Text 1").text_frame.text == "Use cases by department"


def test_apply_slots_handles_list_slots(template_path):
    prs = Presentation(str(template_path))
    new, _ = clone_slide(prs, prs.slides[0])  # cover
    slots = {"steps": ["Text 8", "Text 11", "Text 14", "Text 17", "Text 20"]}
    apply_slots(new, slots, {"steps": ["A", "B", "C", "D", "E"]})
    for name, expected in zip(["Text 8", "Text 11", "Text 14", "Text 17", "Text 20"],
                              ["A", "B", "C", "D", "E"]):
        assert shape_by_name(new, name).text_frame.text == expected


def test_apply_slots_reports_missing(template_path):
    prs = Presentation(str(template_path))
    new, _ = clone_slide(prs, prs.slides[1])
    missing = apply_slots(new, {"title": "Text 1"}, {})
    assert missing == ["title"]
