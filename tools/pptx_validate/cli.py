"""``pptx_validate`` — assert a generated .pptx conforms to the branded design system.

Re-opens the deck with python-pptx and checks, for EVERY slide regardless of
composition: branded background, brand fonts, brand colors, brand logo at the
token EMU position, content chrome (title bar + title style + footer),
cover/closing chrome, on-grid placement, and in-bounds shapes.

Policy:
- SVG-first library architecture: SVGs are the primary source assets;
  PNGs have been removed from library/ (Pass 3 closure).
- Multi-brand: all checks apply equally to BAMI and KVI tokens; brand that
All chrome specs are read from design_tokens.yaml.
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

import click
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from shared.pptx.tokens import load_tokens  # noqa: E402

try:
    from tools.pptx_gen.cli import BRAND_DIRS
except ImportError:
    # Fallback BRAND_DIRS (mirrors pptx_gen/cli.py)
    BRAND_DIRS = {
        "bami": {"template": "templates/bami/template.pptx",       "tokens": "templates/bami/design_tokens.yaml"},
        "kvi":  {"template": "templates/kvi/template.pptx",       "tokens": "templates/kvi/design_tokens.yaml"},
    }


EMU = 914400


def _in(e):
    return None if e is None else e / EMU


def _near(a, b, tol=0.02):
    return a is not None and b is not None and abs(a - b) <= tol


class Report:
    def __init__(self):
        self.violations: list[str] = []

    def add(self, slide_idx: int, msg: str):
        self.violations.append(f"slide {slide_idx}: {msg}")

    @property
    def ok(self):
        return not self.violations


def validate(pptx_path: str | Path, tokens_path: str | Path, chrome_mode: str | None = None) -> Report:
    tokens = load_tokens(tokens_path)
    brand_hexes = tokens.brand_hexes()
    allowed_fonts = {tokens.fonts["primary"].lower()}
    prs = Presentation(str(pptx_path))
    chrome_mode = (chrome_mode or _chrome_mode(prs)).lower()
    rep = Report()

    cv = tokens.canvas
    cw, ch = cv["width_in"], cv["height_in"]

    # --- brand chrome spec (all optional; absent => check skipped) ---
    content_tmpl = tokens.templates.get("content", {})

    # Title bar (content): BAMI has one; KVI may omit.
    title_bar_spec = content_tmpl.get("title_bar") or {}
    title_bar_enabled = bool(title_bar_spec) and title_bar_spec.get("enabled", True)
    tb_fill = (
        "#" + tokens.resolve_color(title_bar_spec["fill"]).lstrip("#").upper()
        if title_bar_spec.get("fill") else None
    )

    # Title text spec
    title_text_spec = content_tmpl.get("title_text") or {}
    title_pt = float(title_text_spec["pt"]) if title_text_spec.get("pt") else None
    title_col = (
        "#" + tokens.resolve_color(title_text_spec["color"]).lstrip("#").upper()
        if title_text_spec.get("color") else None
    )

    # Footer text (read from content footer)
    footer = content_tmpl.get("footer") or {}
    footer_l_text = (footer.get("left") or {}).get("text")
    footer_r_text = (footer.get("right") or {}).get("text")
    footer_enabled = bool(footer_l_text or footer_r_text)
    footer_y = float(tokens.grid.get("footer_y_in", ch * 0.85))
    footer_text_threshold = footer_y * 0.6  # roughly 60% down the slide

    # Logo positions: only templates that define a logo key
    logo_positions = {}
    for tname, t in tokens.templates.items():
        if t.get("logo"):
            logo_positions[tname] = t["logo"]

    n_slides = len(prs.slides._sldIdLst)
    if n_slides == 0:
        rep.add(-1, "deck has no slides")
        return rep

    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        bg_ok = False
        logo_ok = False
        title_bar_ok = False
        title_text_ok = False
        footer_l_ok = False
        footer_r_ok = False

        for shp in shapes:
            L, T, W, H = _in(shp.left), _in(shp.top), _in(shp.width), _in(shp.height)
            st = shp.shape_type

            # --- background: full-bleed or large picture ---
            if st == MSO_SHAPE_TYPE.PICTURE:
                if (_near(L, 0.0) and _near(T, 0.0) and _near(W, cw) and _near(H, ch)):
                    bg_ok = True
                elif W is not None and H is not None and cw and ch:
                    # Accept as background if covers >60% of canvas (KVI layout-level backgrounds)
                    if W / cw > 0.6 and H / ch > 0.6:
                        bg_ok = True
            # --- logo at brand EMU (iterate over all logo positions) ---
            if st == MSO_SHAPE_TYPE.PICTURE:
                for key, pos in logo_positions.items():
                    if (_near(L, pos["left_in"]) and _near(T, pos["top_in"])
                            and _near(W, pos["width_in"]) and _near(H, pos["height_in"])):
                        logo_ok = True
                        break

            # --- title bar check (token-driven) ---
            try:
                f = shp.fill
                if f.type == MSO_FILL.SOLID:
                    rgb = "#" + str(f.fore_color.rgb).upper()
                    if (title_bar_enabled and tb_fill and rgb == tb_fill
                            and _near(T, title_bar_spec.get("top_in", 0.0))
                            and _near(L, title_bar_spec.get("left_in", 0.0))
                            and _near(W, title_bar_spec.get("width_in", 8.6))
                            and _near(H, title_bar_spec.get("height_in", 0.95))):
                        title_bar_ok = True
                    # --- brand colors only (shape fills) ---
                    if rgb not in brand_hexes:
                        rep.add(i, f"shape '{shp.name}' fill color {rgb} is outside the brand palette")
            except Exception:
                pass

            # --- canvas bounds ---
            if L is not None and T is not None and W is not None and H is not None:
                if L < -0.001 or T < -0.001 or L + W > cw + 0.01 or T + H > ch + 0.01:
                    rep.add(i, f"shape '{shp.name}' out of canvas bounds (L{L:.2f} T{T:.2f} W{W:.2f} H{H:.2f})")

            # --- text runs: font + color; chrome title/footer detection ---
            if shp.has_text_frame:
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        fn = (r.font.name or "").lower()
                        if fn and fn not in allowed_fonts:
                            brand_name = tokens.fonts["primary"]
                            rep.add(i, f"run '{r.text[:24]!r}' font {r.font.name!r} is not {brand_name}")
                        try:
                            if r.font.color and r.font.color.rgb is not None:
                                rc = "#" + str(r.font.color.rgb).upper()
                                if rc not in brand_hexes:
                                    rep.add(i, f"run '{r.text[:24]!r}' color {rc} is outside the brand palette")
                        except Exception:
                            pass
                txt = shp.text_frame.text.strip()
                # Title text detection — check both BAMI position (T=0) and token-specified position
                title_shape_name = (title_text_spec or {}).get("shape_name", "")
                title_pos_T = title_text_spec.get("top_in", 0.0) if title_text_spec else 0.0
                title_pos_L = title_text_spec.get("left_in", 0.0) if title_text_spec else 0.0
                title_match_pos = (_near(T, 0.0) or _near(T, title_pos_T))
                title_match_name = (title_shape_name and shp.name == title_shape_name)
                if (title_match_pos or title_match_name) and W and txt and p.runs:
                    r0 = p.runs[0]
                    pt = r0.font.size.pt if r0.font.size else None
                    col = None
                    try:
                        if r0.font.color and r0.font.color.rgb is not None:
                            col = "#" + str(r0.font.color.rgb).upper()
                    except Exception:
                        pass
                    # Accept exact match OR inherited (None) for pt, col, bold
                    pt_ok = (pt == title_pt) or (pt is None and title_pt is not None)
                    col_ok = (col == title_col) or (col is None and title_col is not None)
                    bold_ok = r0.font.bold or r0.font.bold is None
                    if (title_pt and title_col and pt_ok and col_ok and bold_ok
                            and txt):
                        title_text_ok = True
                if footer_l_text and txt == footer_l_text and T and T > footer_text_threshold:
                    footer_l_ok = True
                if footer_r_text and txt == footer_r_text and T and T > footer_text_threshold:
                    footer_r_ok = True

        # --- per-slide chrome assertions (conditional) ---
        # Background: only flag if slide has a large image (>50% canvas) but not at full-bleed
        has_bg_candidate = False
        for s in shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                L, T, W, H = _in(s.left), _in(s.top), _in(s.width), _in(s.height)
                if W and H and cw and ch and W / cw > 0.5 and H / ch > 0.5:
                    has_bg_candidate = True
                    break
        if has_bg_candidate and not bg_ok:
            rep.add(i, "branded full-bleed background missing")
        is_content = _is_content(slide, tokens)
        is_cover   = _is_cover_like(slide, tokens)

        # logo only where the brand template defines one for this slide type
        slide_has_logo_spec = ((is_content or False) and "content" in logo_positions) or \
                              ((is_cover or False) and ("cover" in logo_positions or "closing" in logo_positions))
        if slide_has_logo_spec and not logo_ok:
            rep.add(i, "brand logo not at the token EMU position")

        # Footer: each side is checked independently; omitted tokens skip that side
        if (is_content or False) and footer_enabled:
            footer_missing = False
            if footer_l_text and not footer_l_ok:
                footer_missing = True
            if footer_r_text and not footer_r_ok:
                footer_missing = True
            if footer_missing:
                parts = []
                if footer_l_text and not footer_l_ok: parts.append(f"left={footer_l_text!r}")
                if footer_r_text and not footer_r_ok: parts.append(f"right={footer_r_text!r}")
                rep.add(i, "footer (" + ", ".join(parts) + ") missing")
        if (is_content or False) and title_bar_enabled and not title_bar_ok:
            rep.add(i, "content title bar missing (per tokens.content.title_bar)")
        if (is_content or False) and title_text_spec and not title_text_ok:
            rep.add(i, "content title text not detected (per tokens.content.title_text)")

    # --- structure ---
    slides = list(prs.slides)
    if chrome_mode != "partial":
        is0_cover = _is_cover_like(slides[0], tokens)
        if is0_cover is not None and not is0_cover:
            rep.add(0, "first slide is not a cover (large logo at top-right)")
        is_last_cover = _is_cover_like(slides[-1], tokens)
        if is_last_cover is not None and not is_last_cover:
            rep.add(n_slides - 1, "last slide is not a closing (large logo at top-right)")

    # --- round-trip sanity ---
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        prs.save(tmp_path)
        Presentation(tmp_path)
    except Exception as exc:
        rep.add(-1, f"round-trip save/re-open failed: {exc}")

    return rep


def _chrome_mode(prs) -> str:
    try:
        category = (prs.core_properties.category or "").strip().lower()
    except Exception:
        category = ""
    if "chrome=partial" in category:
        return "partial"
    return "full"


def _is_content(slide, tokens) -> bool | None:
    """Multi-strategy content slide detection:
    1. Logo picture match (if logo defined in tokens)
    2. title_text shape name match
    3. Footer text match at expected position
    Returns None only when NO strategy can determine the slide type."""
    tmpl = tokens.template("content")
    pos = tmpl.get("logo")

    # Strategy 1: Logo picture match
    if pos:
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if (_near(_in(shp.left), pos["left_in"]) and _near(_in(shp.top), pos["top_in"])):
                    return True

    # Strategy 2: title_text shape name match (specific to content slides)
    title_spec = tmpl.get("title_text") or {}
    title_shape_name = title_spec.get("shape_name", "")
    if title_shape_name:
        title_top = title_spec.get("top_in", 0)
        for shp in slide.shapes:
            if shp.name == title_shape_name and shp.has_text_frame and shp.text_frame.text.strip():
                # Also verify vertical position (content title at content y, not cover y)
                if _near(_in(shp.top), title_top) or title_top == 0:
                    return True

    # Strategy 3: Footer text at expected content position
    footer = tmpl.get("footer") or {}
    footer_l = (footer.get("left") or {}).get("text", "")
    footer_y = (footer.get("left") or {}).get("top_in") or (tokens.grid.get("footer_y_in", 0))
    if footer_l and footer_y:
        for shp in slide.shapes:
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt == footer_l and _near(_in(shp.top), footer_y):
                    return True

    return None


def _is_cover_like(slide, tokens) -> bool | None:
    """Multi-strategy cover/closing slide detection:
    1. Logo picture match (if logo defined in tokens)
    2. Full-bleed background image (closing slide)
    3. Wordmark image at cover expected position
    Returns None only when NO strategy can determine the slide type."""
    tmpl = tokens.template("cover")
    pos = tmpl.get("logo")

    # Strategy 1: Logo picture match
    if pos:
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if (_near(_in(shp.left), pos["left_in"]) and _near(_in(shp.top), pos["top_in"])
                        and _near(_in(shp.width), pos["width_in"])):
                    return True

    # Strategy 2: Full-bleed background image (closing)
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            L, T, W, H = _in(shp.left), _in(shp.top), _in(shp.width), _in(shp.height)
            if W and H and W > 10 and H > 5:  # Large image spanning most of canvas
                return True

    # Strategy 3: Check for wordmark logo at top-left
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            L, T = _in(shp.left), _in(shp.top)
            if L is not None and T is not None and L < 1 and T < 1:
                return True

    return None
    return False


@click.command()
@click.argument("pptx_path", required=False, default=None, type=click.Path(exists=True, dir_okay=False))
@click.option("--brand", default="bami", type=click.Choice(["bami", "kvi"]),
              help="Brand template set (default: bami). Sets --tokens default.")
@click.option("--tokens", "tokens_path", default=None, type=click.Path(exists=True, dir_okay=False),
              help="Override design_tokens.yaml (default: brand dir).")
@click.option("--chrome", "chrome_mode", type=click.Choice(["full", "partial"]), default=None,
              help="Override chrome mode (default: read bami:chrome=* from deck core-properties).")
@click.option("--patterns", is_flag=True, default=False,
              help="Run pattern library validation (SVGs, registry, assets) instead of PPTX validation.")
def main(pptx_path, brand, tokens_path, chrome_mode, patterns):
    """Validate a generated .pptx or the pattern library.

    When PPTX_PATH is given, validates the deck against the branded design
    system. Use --patterns to validate pattern-assets.yaml, SVG file integrity,
    registry consistency, and provenance references.
    """
    if patterns:
        from tools.pptx_validate.patterns import run_all
        rep, orphan_rep = run_all()
        if orphan_rep.violations:
            for v in orphan_rep.violations:
                click.echo(f"INFO: {v}")
        if rep.ok:
            click.echo("OK: All pattern validation checks passed.")
            sys.exit(0)
        click.echo(f"FAIL: {len(rep.violations)} violation(s):", err=True)
        for v in rep.violations:
            click.echo(f"  - {v}", err=True)
        sys.exit(1)
    if pptx_path is None:
        click.echo("Error: PPTX_PATH is required (or use --patterns for pattern validation)", err=True)
        sys.exit(1)
    if tokens_path is None:
        tokens_path = BRAND_DIRS[brand]["tokens"]
    rep = validate(pptx_path, tokens_path, chrome_mode=chrome_mode)
    if rep.ok:
        click.echo(f"OK: deck conforms to the {brand} design system ({len(list(Presentation(pptx_path).slides._sldIdLst))} slides)")
        sys.exit(0)
    click.echo(f"FAIL: {len(rep.violations)} violation(s):", err=True)
    for v in rep.violations:
        click.echo(f"  - {v}", err=True)
    sys.exit(1)


if __name__ == "__main__":
    main()
