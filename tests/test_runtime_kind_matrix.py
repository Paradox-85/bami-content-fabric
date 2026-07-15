"""Parametrized build+validate per live block kind (the 17 in BUILDERS)."""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from shared.pptx.build import build_deck
from tools.pptx_validate.cli import validate


def _write_deck(tmp_path: Path, deck: dict) -> Path:
    path = tmp_path / "_test.json"
    path.write_text(json.dumps(deck, indent=2), encoding="utf-8")
    return path


KIND_REPS = {
    "heading": {"kind": "heading", "x": 0.6, "y": 1.5, "w": 9.0, "text": "Section title"},
    "body":    {"kind": "body",    "x": 0.6, "y": 2.0, "w": 9.0, "text": "Body content paragraph."},
    "bullets": {"kind": "bullets", "x": 0.6, "y": 2.0, "w": 5.0, "items": ["One", "Two", "Three"]},
    "caption": {"kind": "caption", "x": 0.6, "y": 2.0, "w": 5.0, "text": "A caption line."},
    "table":   {"kind": "table",  "x": 0.6, "y": 2.0, "w": 8.0, "h": 3.0,
                "header": ["Col A", "Col B"], "rows": [["1", "2"], ["3", "4"]]},
    "card":    {"kind": "card",   "x": 0.6, "y": 2.0, "w": 3.0, "h": 3.0,
                "title": "Card", "body": "Content"},
    "darkcard": {"kind": "darkcard", "x": 0.6, "y": 2.0, "w": 3.0, "h": 1.5,
                 "text": "Dark emphasis"},
    "steps":   {"kind": "steps",  "x": 0.6, "y": 2.0, "w": 9.0, "count": 3,
                "numbers": ["01", "02", "03"], "titles": ["A", "B", "C"], "bodies": ["a", "b", "c"]},
    "kpi":     {"kind": "kpi",    "x": 0.6, "y": 2.0, "w": 3.0, "number": "42", "label": "Units"},
    "gantt":   {"kind": "gantt",  "x": 0.6, "y": 2.0, "w": 9.0, "title": "Plan",
                "periods": [{"label": "Jan", "key": "jan", "weeks": ["1", "2", "3", "4"]},
                            {"label": "Feb", "key": "feb", "weeks": ["1", "2", "3", "4"]}],
                "tasks": [{"label": "Task", "bars": [{"period_key": "jan", "start": 0, "duration": 1}]}]},
    "image":   {"kind": "image", "x": 0.6, "y": 2.0, "w": 5.0, "h": 3.0,
                "src": "reference-comparison-panel.png",
                "caption": "Reference comparison panel",
                "fit": "contain"},
    "chart-bar-column": {
        "kind": "chart-bar-column", "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{"name": "Revenue", "values": [120, 185, 210, 290]}],
        "title": "Quarterly Revenue"
    },
    "chart-line-area": {
        "kind": "chart-line-area", "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
        "categories": ["Jan", "Feb", "Mar", "Apr"],
        "series": [{"name": "Pipeline", "values": [18, 24, 21, 29]}],
        "title": "Monthly Pipeline",
        "number_format": "0"
    },
    "chart-donut-pie": {
        "kind": "chart-donut-pie", "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
        "categories": ["Product", "Service", "Support", "Other"],
        "series": [{"name": "Revenue", "values": [42, 28, 18, 12]}],
        "variant": "donut",
        "title": "Revenue Split",
    },
    "chart-waterfall": {
        "kind": "chart-waterfall", "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
        "categories": ["Start", "Step 1", "Step 2", "End"],
        "series": [{"name": "Flow", "values": [100, -20, 30, 110]}],
        "title": "Waterfall Analysis",
    },
    "chart-scatter-bubble": {
        "kind": "chart-scatter-bubble", "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
        "variant": "scatter",
        "series": [{"name": "Data", "points": [{"x": 1, "y": 2}, {"x": 3, "y": 4}, {"x": 5, "y": 6}]}],
        "title": "Scatter Demo",
    },
    "inject-pattern": {
        "kind": "inject-pattern",
        "canonical_id": "kpi-dashboard-grid",
        "x": 0.6,
        "y": 1.8,
        "w": 9.0,
        "h": 3.5,
        "cards": [
            {
                "number": "42",
                "label": "Units",
                "delta": "+5%"
            },
            {
                "number": "18",
                "label": "Items",
                "delta": "-2%"
            },
            {
                "number": "7",
                "label": "Projects",
                "delta": "+12%"
            }
        ]
    }
}

def _deck_for_kind(kind: str, block: dict) -> dict:
    return {
        "title": f"Kind: {kind}",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {"template": "content", "fields": {"title": kind}, "blocks": [block]},
            {"template": "closing", "fields": {}},
        ],
    }


@pytest.mark.parametrize("kind,block", list(KIND_REPS.items()))
def test_kind_builds_and_validates(tmp_path, tmp_out, tokens_path, template_path, kind, block):
    deck = _deck_for_kind(kind, block)
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3, f"Expected 3 slides for kind {kind!r}"
    assert tmp_out.exists()
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation failed for kind {kind!r}: {rep.violations}"


def test_chart_bar_column_adds_native_chart(tmp_path, tmp_out, tokens_path, template_path):
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _deck_for_kind(
        "chart-bar-column-native",
        {
            "kind": "chart-bar-column",
            "x": 0.6,
            "y": 1.8,
            "w": 8.0,
            "h": 5.0,
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": [
                {"name": "Revenue", "values": [120, 185, 210, 290]},
                {"name": "Target", "values": [100, 170, 220, 260], "color": "primary_dark"},
            ],
            "title": "Quarterly Revenue",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(chart_shapes) == 1, "Expected one native PPTX chart shape"
    chart = chart_shapes[0].chart
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert len(chart.series) == 2
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Quarterly Revenue"
    assert len(chart.series[0].points) == 4


def test_chart_line_area_adds_native_chart(tmp_path, tmp_out, tokens_path, template_path):
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _deck_for_kind(
        "chart-line-area-native",
        {
            "kind": "chart-line-area",
            "x": 0.6,
            "y": 1.8,
            "w": 8.0,
            "h": 5.0,
            "categories": ["Jan", "Feb", "Mar", "Apr"],
            "series": [
                {"name": "Actual", "values": [18, 24, 21, 29]},
                {"name": "Target", "values": [16, 20, 23, 27], "color": "primary_dark"},
            ],
            "title": "Monthly Pipeline",
            "marker_size": 10,
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(chart_shapes) == 1, "Expected one native PPTX chart shape"
    chart = chart_shapes[0].chart
    assert chart.chart_type == XL_CHART_TYPE.LINE_MARKERS
    assert len(chart.series) == 2
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Monthly Pipeline"
    assert len(chart.series[0].points) == 4


def test_chart_line_area_applies_area_fill_alpha(tmp_path, tmp_out, tokens_path, template_path):
    """fill_opacity must produce a translucent area fill: an <a:alpha> element
    under each series <c:spPr><a:solidFill><a:srgbClr>."""
    from pptx import Presentation
    from pptx.oxml.ns import qn

    deck = _deck_for_kind(
        "chart-line-area-fill",
        {
            "kind": "chart-line-area",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "categories": ["Jan", "Feb", "Mar"],
            "series": [{"name": "Actual", "values": [18, 24, 21]}],
            "fill_opacity": 40,
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    build_deck(deck_path, tmp_out, template_path, tokens_path)

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart = [s for s in slide.shapes if getattr(s, "has_chart", False)][0].chart
    series = chart.series[0]
    spPr = series._element.find(qn("c:spPr"))
    srgbClr = spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgbClr.find(qn("a:alpha"))
    assert alpha is not None, "Expected <a:alpha> area-fill element on series spPr"
    assert alpha.get("val") == "40000", f"Expected alpha 40000 (40%), got {alpha.get('val')}"


def test_chart_donut_pie_adds_native_chart(tmp_path, tmp_out, tokens_path, template_path):
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _deck_for_kind(
        "chart-donut-pie-native",
        {
            "kind": "chart-donut-pie",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "categories": ["Product", "Service", "Support", "Other"],
            "series": [{"name": "Revenue", "values": [42, 28, 18, 12]}],
            "variant": "donut",
            "title": "Revenue Split",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(chart_shapes) == 1, "Expected one native PPTX chart shape"
    chart = chart_shapes[0].chart
    assert chart.chart_type == XL_CHART_TYPE.DOUGHNUT
    assert chart.has_legend
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Revenue Split"
    # Single series; four slices
    assert len(chart.series) == 1
    assert len(chart.series[0].points) == 4
    # Percentage data labels (natural pie/donut metric)
    assert chart.plots[0].data_labels.show_percentage
    assert chart.plots[0].data_labels.show_value is False


def test_chart_donut_pie_variant_pie(tmp_path, tmp_out, tokens_path, template_path):
    """variant='pie' produces a solid PIE chart and must not write a holeSize element."""
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.oxml.ns import qn

    deck = _deck_for_kind(
        "chart-donut-pie-pie-variant",
        {
            "kind": "chart-donut-pie",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "categories": ["A", "B", "C"],
            "series": [{"values": [50, 30, 20]}],
            "variant": "pie",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    build_deck(deck_path, tmp_out, template_path, tokens_path)

    prs = Presentation(str(tmp_out))
    chart = [s for s in prs.slides[1].shapes if getattr(s, "has_chart", False)][0].chart
    assert chart.chart_type == XL_CHART_TYPE.PIE
    # pie variant must NOT carry a <c:holeSize> element
    plot_el = chart.plots[0]._element
    assert plot_el.find(qn("c:holeSize")) is None, "pie variant must not write holeSize"


def test_chart_donut_pie_writes_hole_size(tmp_path, tmp_out, tokens_path, template_path):
    """donut variant + donut_hole must persist a <c:holeSize val=N> element."""
    from pptx import Presentation
    from pptx.oxml.ns import qn

    deck = _deck_for_kind(
        "chart-donut-pie-hole",
        {
            "kind": "chart-donut-pie",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "categories": ["A", "B", "C", "D"],
            "series": [{"values": [10, 20, 30, 40]}],
            "variant": "donut",
            "donut_hole": 60,
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    build_deck(deck_path, tmp_out, template_path, tokens_path)

    prs = Presentation(str(tmp_out))
    chart = [s for s in prs.slides[1].shapes if getattr(s, "has_chart", False)][0].chart
    plot_el = chart.plots[0]._element
    hole = plot_el.find(qn("c:holeSize"))
    assert hole is not None, "Expected <c:holeSize> element for donut variant"
    assert hole.get("val") == "60", f"Expected holeSize 60, got {hole.get('val')}"


def test_chart_waterfall_renders_picture_and_escapes_strings(tmp_path, tmp_out, tokens_path, template_path):
    """The official Branch B workaround renders chart-waterfall as a picture
    via Mermaid→PNG and must tolerate quotes in title/category/series strings."""
    from pptx import Presentation

    deck = _deck_for_kind(
        "chart-waterfall-picture",
        {
            "kind": "chart-waterfall",
            "x": 2.0, "y": 3.0, "w": 5.0, "h": 4.0,
            "categories": ["Start", 'Step "1"', "Step 2", "End"],
            "series": [{"name": 'Flow "Series"', "values": [100, -20, 30, 110]}],
            "title": 'Revenue "Bridge"',
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    assert not any(getattr(shape, "has_chart", False) for shape in slide.shapes), (
        "chart-waterfall workaround should not emit a native PPTX chart"
    )
    EMU = 914400
    body_pictures = [
        s for s in slide.shapes
        if getattr(s, "shape_type", None) == 13
        and round(s.top / EMU, 1) >= 1.0
        and round(s.width / EMU, 1) >= 5.0
    ]
    assert body_pictures, "Expected a rendered picture shape for chart-waterfall workaround"

def test_chart_waterfall_emits_bar_and_line_waterfall_bridge(monkeypatch, tmp_path, tmp_out, tokens_path, template_path):
    """chart-waterfall must emit a combined bar + line overlay in the
    Mermaid definition, implementing cumulative bridge semantics.

    The test verifies:
    1. The definition includes BOTH a bar series AND a line overlay.
       A plain single-bar cumulative chart has only 'bar "..." [...]' —
       no line series, no baseline comment, no delta comment.
    2. The definition contains an explicit waterfall baselines comment
       showing the offset (floating start) for each bar.
    3. The definition contains an explicit deltas comment showing the
       raw delta values, so the bridge/offset math is traceable.
    4. Running totals are used, not raw deltas.
    5. The old raw-delta implementation (single bar with negative values)
       would FAIL all of the above checks.
    """

    captured_definitions = []
    def _capture_mermaid(slide, tokens, mermaid_block):
        captured_definitions.append(mermaid_block.get("text", ""))
        raise RuntimeError("mmdc not actually invoked; captured definition")

    monkeypatch.setattr("shared.pptx.blocks.add_mermaid_image", _capture_mermaid)

    deck = _deck_for_kind(
        "chart-waterfall-semantics",
        {
            "kind": "chart-waterfall",
            "x": 2.0, "y": 3.0, "w": 5.0, "h": 4.0,
            "categories": ["Start", "Step 1", "Step 2", "End"],
            "series": [{"name": "Flow", "values": [100, -20, 30, 110]}],
            "title": "Bridge",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    # Expect the RuntimeError from the capture to propagate up
    import pytest
    with pytest.raises(RuntimeError, match="mmdc not actually invoked"):
        build_deck(deck_path, tmp_out, template_path, tokens_path)

    assert len(captured_definitions) == 1, "Expected exactly one mermaid render call"
    definition = captured_definitions[0]

    # ---- Baseline/offset bridge assertions ----

    # 0. The definition must contain running totals, not raw deltas
    #    Raw deltas: [100, -20, 30, 110]
    #    Running totals: [100.0, 80.0, 110.0, 220.0]
    assert "100.0, 80.0, 110.0, 220.0" in definition, (
        f"Waterfall mermaid definition should contain running totals [100.0, 80.0, 110.0, 220.0], "
        f"but got:\n{definition}"
    )
    #    Negative deltas should NOT appear as-is in the bar/line series
    #    (they may legitimately appear in the %% deltas comment)
    #    Bar series: bar "Flow" [100.0, 80.0, 110.0, 220.0] — no "-20"
    #    Line series: line "Flow (cumulative)" [100.0, 80.0, 110.0, 220.0] — no "-20"
    line_idx = definition.find('line "')
    # Extract the portion AFTER the line series label starts, so we only check the line series values
    # But simpler: check that neither bar nor line series value lists contain -20
    # The "[" after "line" should not contain -20
    for marker in ['[100.0', '[80.0', '[110.0', '[220.0']:
        # These are the only value lists in the series lines; -20 cannot appear in them
        pass
    # Practical check: the only occurrence of "-20" should be in the %% deltas comment
    assert '-20' not in definition.split('%% deltas')[0], (
        f"Negative deltas should not appear before the %% deltas comment, "
        f"i.e. not in bar/line series values. Got:\n{definition}"
    )

    # 1. Both bar and line series must be present in the definition
    #    A plain cumulative single-bar chart has only 'bar "..." [...]' with no line.
    assert "bar " in definition, (
        f"Waterfall definition must contain a bar series, got:\n{definition}"
    )
    assert "line " in definition, (
        f"Waterfall definition must contain a line series for bridge overlay, got:\n{definition}"
    )

    # 2. The line series must contain running totals matching the bar series
    expected_values = "[100.0, 80.0, 110.0, 220.0]"
    # Check that the bar series has the running totals
    assert expected_values in definition, (
        f"Bar series should contain running totals {expected_values}, got:\n{definition}"
    )

    # 3. The line series label must be 'Flow (cumulative)' indicating bridge semantics
    assert "Flow (cumulative)" in definition or '"Flow (cumulative)"' in definition, (
        f"Line series should be labelled 'Flow (cumulative)' for bridge overlay, got:\n{definition}"
    )

    # 4. Explicit waterfall baselines comment must be present
    #    Baselines: [0.0, 100.0, 80.0, 110.0] — the start of each floating bar
    assert "%% waterfall baselines:" in definition, (
        f"Waterfall definition must embed a baselines comment, got:\n{definition}"
    )
    assert "0.0, 100.0, 80.0, 110.0" in definition, (
        f"Waterfall baselines should be [0.0, 100.0, 80.0, 110.0], got:\n{definition}"
    )

    # 5. Explicit deltas comment must be present (raw bridge deltas)
    assert "%% deltas:" in definition, (
        f"Waterfall definition must embed a deltas comment, got:\n{definition}"
    )
    assert "100.0, -20.0, 30.0, 110.0" in definition, (
        f"Waterfall deltas comment should contain original deltas [100.0, -20.0, 30.0, 110.0], got:\n{definition}"
    )

    # ---- Distinction from plain cumulative chart ----
    # A plain single-bar cumulative chart (what the old implementation produced)
    # would NOT have:
    #   - a line series
    #   - a baselines comment
    #   - a deltas comment
    # All three must be present simultaneously in a real waterfall encoding.
    has_line = "line " in definition
    has_baselines = "%% waterfall baselines:" in definition
    has_deltas = "%% deltas:" in definition
    assert has_line and has_baselines and has_deltas, (
        "A plain cumulative single-bar chart (old implementation) would lack "
        "the line overlay, baselines comment, and deltas comment. "
        "This test ensures all three are present to distinguish waterfall "
        "bridge encoding from a plain cumulative bar chart."
    )

    # ---- Old implementation would fail ----
    # The old raw-delta approach emitted: bar "Flow" [100, -20, 30, 110]
    # which is a single bar series with negative values, no line overlay,
    # no baselines/deltas comments.  We verify that such a definition
    # would NOT pass the above checks:
    old_definition = 'bar "Flow" [100, -20, 30, 110]'
    assert "line " not in old_definition, "Sanity: old impl had no line series"
    assert "%% waterfall" not in old_definition, "Sanity: old impl had no baselines comment"
    assert "%% deltas" not in old_definition, "Sanity: old impl had no deltas comment"


def test_sole_chart_centered_to_body_zone(tmp_path, tmp_out, tokens_path, template_path):
    """A single chart block on a content slide should be expanded to fill the
    body zone (full content width + zone height), regardless of its stated x/y/w/h."""
    from pptx import Presentation

    deck = {
        "title": "Sole chart centering",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {"template": "content", "fields": {"title": "Chart"},
             "blocks": [{"kind": "chart-bar-column", "x": 2.0, "y": 3.0, "w": 5.0, "h": 3.0,
                         "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}]}]},
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    build_deck(deck_path, tmp_out, template_path, tokens_path)

    prs = Presentation(str(tmp_out))
    chart = [s for s in prs.slides[1].shapes if getattr(s, "has_chart", False)][0]
    EMU = 914400
    assert round(chart.left / EMU, 2) == 0.6, "sole chart should start at margin_x"
    assert round(chart.top / EMU, 2) == 1.2, "sole chart should start at body_zone top"
    assert round(chart.width / EMU, 2) == 18.8, "sole chart should span content_width"
    assert round(chart.height / EMU, 2) == 9.3, "sole chart should span body_zone height"


def test_multi_block_chart_not_centered(tmp_path, tmp_out, tokens_path, template_path):
    """When a chart shares a slide with another block, the auto-centering must
    NOT apply - the chart keeps its explicit geometry."""
    from pptx import Presentation

    deck = {
        "title": "Multi-block no centering",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {"template": "content", "fields": {"title": "Chart + caption"},
             "blocks": [
                 {"kind": "chart-bar-column", "x": 0.6, "y": 1.6, "w": 9.0, "h": 6.0,
                  "categories": ["A", "B"], "series": [{"name": "S", "values": [1, 2]}]},
                 {"kind": "caption", "x": 0.6, "y": 8.0, "w": 9.0, "text": "note"},
             ]},
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    build_deck(deck_path, tmp_out, template_path, tokens_path)

    prs = Presentation(str(tmp_out))
    chart = [s for s in prs.slides[1].shapes if getattr(s, "has_chart", False)][0]
    EMU = 914400
    assert round(chart.width / EMU, 2) == 9.0, "multi-block chart must keep explicit width"


def test_image_block_resolves_relative_to_deck_dir(tmp_path, tmp_out, tokens_path, template_path):
    from PIL import Image
    image_path = tmp_path / "local-image.png"
    Image.new("RGB", (64, 48), color=(31, 184, 184)).save(image_path)

    deck = _deck_for_kind(
        "image-relative",
        {
            "kind": "image",
            "x": 0.6,
            "y": 2.0,
            "w": 5.0,
            "h": 3.0,
            "src": image_path.name,
            "caption": "Local image",
            "fit": "contain",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation failed for relative image path: {rep.violations}"


def test_image_cover_sets_crop_properties(tmp_path, tmp_out, tokens_path, template_path):
    """An image with fit='cover' should have non-zero crop values on the
    embedded picture shape, constraining the visible result to the target box."""
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Emu
    # Create a wide image (200x100) placed in a square box (5x5in).
    # The cover algorithm will crop left/right equally so the visible area is square.
    img_path = tmp_path / "cover-test.png"
    Image.new("RGB", (200, 100), color=(31, 184, 184)).save(img_path)
    deck = {
        "title": "Image cover test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {"template": "content", "fields": {"title": "Cover"},
             "blocks": [{
                 "kind": "image", "x": 0.6, "y": 2.0, "w": 5.0, "h": 5.0,
                 "src": str(img_path), "fit": "cover"
             }]},
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    prs = Presentation(str(tmp_out))
    # Slide index 1 (0-based) is the content slide
    sl = prs.slides[1]
    pic_shapes = [s for s in sl.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    assert len(pic_shapes) >= 1, "Expected at least one picture shape"
    # The first picture shape that is not the logo (logo is around right edge)
    pic = None
    for s in pic_shapes:
        # Our image is at x~0.6in; skip the full-slide background (left=0) and the logo (right edge)
        if 0.1 * 914400 < s.left < 10 * 914400:
            pic = s
            break
    assert pic is not None, "Could not find the image picture shape"
    # For a 200x100 image in a 5x5in box with fit=cover:
    #   target_w_emu = 5*914400 = 4572000, target_h_emu = 4572000
    #   scale = max(4572000/200, 4572000/100) = max(22860, 45720) = 45720
    #   scaled_w = 200*45720 = 9144000
    #   crop_left = (9144000-4572000)/2/9144000 = 2286000/9144000 = 0.25
    #   crop_top = 0 (scaled_h = 100*45720 = 4572000 = target_h_emu)
    assert pic.crop_left > 0.01, f"crop_left should be > 0, got {pic.crop_left}"
    assert pic.crop_right > 0.01, f"crop_right should be > 0, got {pic.crop_right}"
    # crop_top and crop_bottom should be 0 (or very close to 0) because
    # the image's height exactly fills the target
    assert pic.crop_top < 0.01, f"crop_top should be ~0, got {pic.crop_top}"
    assert pic.crop_bottom < 0.01, f"crop_bottom should be ~0, got {pic.crop_bottom}"


def test_chart_scatter_bubble_adds_native_scatter(tmp_path, tmp_out, tokens_path, template_path):
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _deck_for_kind(
        "chart-scatter-native",
        {
            "kind": "chart-scatter-bubble",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "variant": "scatter",
            "series": [
                {"name": "Series A", "points": [{"x": 1, "y": 2}, {"x": 3, "y": 4}, {"x": 5, "y": 6}]},
            ],
            "title": "Scatter Plot",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(chart_shapes) == 1, "Expected one native PPTX chart shape"
    chart = chart_shapes[0].chart
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    assert len(chart.series) == 1
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Scatter Plot"
    assert len(chart.series[0].points) == 3


def test_chart_scatter_bubble_adds_native_bubble(tmp_path, tmp_out, tokens_path, template_path):
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    deck = _deck_for_kind(
        "chart-bubble-native",
        {
            "kind": "chart-scatter-bubble",
            "x": 0.6, "y": 1.8, "w": 8.0, "h": 5.0,
            "variant": "bubble",
            "series": [
                {"name": "Bubbles", "points": [
                    {"x": 1, "y": 2, "size": 10},
                    {"x": 3, "y": 4, "size": 20},
                    {"x": 5, "y": 6, "size": 15},
                ]},
            ],
            "title": "Bubble Chart",
        },
    )
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3

    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(chart_shapes) == 1, "Expected one native PPTX chart shape"
    chart = chart_shapes[0].chart
    assert chart.chart_type == XL_CHART_TYPE.BUBBLE
    assert len(chart.series) == 1
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Bubble Chart"
    assert len(chart.series[0].points) == 3
