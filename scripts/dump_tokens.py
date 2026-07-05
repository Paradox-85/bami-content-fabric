#!/usr/bin/env python
"""Print the template's chrome geometry so design_tokens.yaml can be verified.

For each of the three reference slides (cover/content/closing) it lists the
picture shapes (background + logo) and the named text shapes, with position,
size and the first run's font/size/color/bold/alignment. Use this after any
change to templates/template.pptx to confirm the slot maps in design_tokens.yaml
still match reality.

Usage (from the repository root; repository identity is transitioning from
presentation-framework to bami-content-fabric):
    python scripts/dump_tokens.py
    python scripts/dump_tokens.py --template templates/template.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DEFAULT_TEMPLATE = "templates/template.pptx"
REFERENCES = {"cover": 0, "content": 1, "closing": 7}


def _in(e):
    return None if e is None else round(e / 914400, 3)


def _run(tf):
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        return None
    r = tf.paragraphs[0].runs[0]
    col = None
    try:
        if r.font.color and r.font.color.rgb is not None:
            col = "#" + str(r.font.color.rgb)
    except Exception:
        pass
    return {
        "text": r.text[:40],
        "font": r.font.name,
        "pt": r.font.size.pt if r.font.size else None,
        "bold": r.font.bold,
        "color": col,
        "align": str(tf.paragraphs[0].alignment),
    }


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", default=DEFAULT_TEMPLATE)
    args = ap.parse_args()

    prs = Presentation(args.template)
    print(f"# dump_tokens: {args.template}  ({len(prs.slides._sldIdLst)} slides)")
    print(f"# canvas: {_in(prs.slide_width)} x {_in(prs.slide_height)} in")
    print()
    for name, idx in REFERENCES.items():
        if idx >= len(prs.slides):
            print(f"[!] {name}: ref_index {idx} out of range — skipping")
            continue
        sl = prs.slides[idx]
        print(f"## {name} (slide {idx}) — {len(sl.shapes)} shapes")
        for shp in sl.shapes:
            st = "PIC" if shp.shape_type == MSO_SHAPE_TYPE.PICTURE else "TXT"
            pos = f"({_in(shp.left)},{_in(shp.top)}) size=({_in(shp.width)},{_in(shp.height)})"
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"  [{st}] {shp.name:<14} {pos}")
            elif shp.has_text_frame and shp.text_frame.text.strip():
                r = _run(shp.text_frame)
                extra = f"  font={r['font']} pt={r['pt']} bold={r['bold']} color={r['color']} align={r['align']}" if r else ""
                print(f"  [{st}] {shp.name:<14} {pos}  '{r['text']}'{extra}" if r else f"  [{st}] {shp.name:<14} {pos}")
        print()


if __name__ == "__main__":
    main()
