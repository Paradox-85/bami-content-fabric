"""Tests for the simple-arrow-horizontal native pattern injector.

Verifies:
- Injector is registered and callable
- B1 regression: the invalid color token (neutral_light) is fixed
- Resolver can select simple-arrow-horizontal via explicit graphical_variant
- Build pipeline works with explicit graphical_variant selection
- Shape naming convention is followed
"""

from __future__ import annotations

import json
from pathlib import Path

from pathlib import Path

from pptx import Presentation
from shared.pptx.pattern_injectors.registry import get_injector, list_injectors, inject_pattern
from shared.pptx.pattern_selection import resolve_pattern
from tools.pptx_validate.cli import validate


# ---------------------------------------------------------------------------
# Unit: injector registration
# ---------------------------------------------------------------------------

def test_simple_arrow_injector_registered():
    """The simple-arrow-horizontal injector is registered and callable."""
    assert "simple-arrow-horizontal" in list_injectors()
    fn = get_injector("simple-arrow-horizontal")
    assert fn is not None
    assert callable(fn)


# ---------------------------------------------------------------------------
# Unit: B1 regression — invalid color token
# ---------------------------------------------------------------------------

def test_simple_arrow_no_invalid_color_token():
    """The injector does not reference the undefined 'neutral_light' token.

    Regression for Blocker B1. We verify by calling the injector with
    >=2 steps (which exercises the connector path that previously crashed).
    """
    from shared.pptx.tokens import load_tokens
    ROOT = Path(__file__).resolve().parent.parent
    real_tokens = load_tokens(ROOT / "templates" / "bami" / "design_tokens.yaml")
    # Create a minimal slide stub suitable for inject_pattern
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Call inject_pattern with 3 steps — this previously crashed at connector rendering
    created = inject_pattern(
        slide, real_tokens,
        "simple-arrow-horizontal",
        x=0.5, y=1.0, w=9.0, h=3.0,
        steps=[
            {"number": "01", "title": "Plan"},
            {"number": "02", "title": "Build"},
            {"number": "03", "title": "Deploy"},
        ],
    )
    assert len(created) >= 5, (
        f"Expected at least 5 shapes (3 circles + 3 numbers + 2 connectors + 3 titles), "
        f"got {len(created)}"
    )


# ---------------------------------------------------------------------------
# Unit: resolver enrichment — explicit graphical_variant
# ---------------------------------------------------------------------------

def test_resolver_returns_simple_arrow_variant():
    """Content with items + explicit graphical_variant -> simple-arrow-horizontal."""
    class FakeTokens:
        def __init__(self, brand="bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    tokens = FakeTokens("bami")
    sel = resolve_pattern(
        content, tokens,
        graphical_variant="simple-arrow-horizontal",
    )
    assert sel.family == "numbered-process-steps"
    assert sel.pattern_template_id == "numbered-process-steps/simple-arrow-horizontal@1.0.0"
    assert sel.graphical_variant == "simple-arrow-horizontal"
    assert sel.renderer_binding is not None
    assert sel.renderer_binding["native"]["injector_id"] == "simple-arrow-horizontal"
    assert sel.features is not None
    assert sel.features.get("native_editable") is True


# ---------------------------------------------------------------------------
# Integration: build pipeline with explicit graphical_variant
# ---------------------------------------------------------------------------

def test_build_pipeline_simple_arrow_with_explicit_graphical_variant(tmp_path, tmp_out, tokens_path, template_path):
    """A deck with graphical_variant: simple-arrow-horizontal resolves and builds."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Simple Arrow Build Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Process"},
                "graphical_variant": "simple-arrow-horizontal",
                "content": {
                    "items": ["Step A", "Step B", "Step C"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_simple_arrow_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: shape naming convention in generated PPTX
# ---------------------------------------------------------------------------

def test_simple_arrow_shape_naming_convention(tmp_path, tmp_out, tokens_path, template_path):
    """The generated PPTX contains shapes with the simple-arrow naming convention."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Simple Arrow Shape Naming Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Named Process"},
                "graphical_variant": "simple-arrow-horizontal",
                "content": {
                    "items": ["Alpha", "Beta", "Gamma"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_simple_arrow_shape_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    # Open the PPTX and inspect shape names on the content slide (index 1)
    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    shape_names = [shp.name for shp in slide.shapes]
    # Check naming convention pattern
    pattern_shapes = [n for n in shape_names if n.startswith("pattern:numbered-process-steps/simple-arrow-horizontal")]
    assert len(pattern_shapes) >= 3, (
        f"Expected at least 3 pattern-named shapes, got {len(pattern_shapes)}: {pattern_shapes}"
    )
    # Verify role naming
    circle_shapes = [n for n in pattern_shapes if n.endswith(":circle")]
    number_shapes = [n for n in pattern_shapes if n.endswith(":number")]
    title_shapes = [n for n in pattern_shapes if n.endswith(":title")]
    connector_shapes = [n for n in pattern_shapes if ":connector:" in n]
    assert len(circle_shapes) == 3, f"Expected 3 circle shapes, got {len(circle_shapes)}"
    assert len(number_shapes) == 3, f"Expected 3 number shapes, got {len(number_shapes)}"
    assert len(title_shapes) == 3, f"Expected 3 title shapes, got {len(title_shapes)}"
    assert len(connector_shapes) == 2, f"Expected 2 connector shapes, got {len(connector_shapes)}"

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# KVI brand variant tests
# ---------------------------------------------------------------------------


def test_simple_arrow_kvi_builds_ok(tmp_path, tmp_out, kvi_template_path, kvi_tokens_path):
    """KVI-branded deck with simple-arrow-horizontal resolves and builds."""
    from shared.pptx.build import build_deck
    deck = {
        "title": "KVI Simple Arrow Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "KVI Steps"},
                "graphical_variant": "simple-arrow-horizontal",
                "content": {
                    "items": ["Step A", "Step B", "Step C"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_kvi_simple_arrow_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")
    result = build_deck(deck_path, tmp_out, kvi_template_path, kvi_tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    rep = validate(tmp_out, kvi_tokens_path)
    assert rep.ok, f"KVI validation violations: {rep.violations}"
