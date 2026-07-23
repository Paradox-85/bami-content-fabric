"""Tests for the visual-fidelity validator (shared/pptx/visual_fidelity.py).

Verifies:
- Fidelity stage gate (semantic-only/placeholder → not enabled)
- Registry-wide fidelity gate scan
- Graphical area sufficiency
- Shape count within budget
- White-on-white detection
- Color roles sufficiency
- Text-to-graphics ratio
- Spatial balance
- Unnatural short wrapping detection
- Required icons/icon count enforcement
- Required visual layers detection
- Full-report aggregation
"""
from __future__ import annotations

from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from shared.pptx import visual_fidelity as vf
from shared.pptx.visual_fidelity import NON_CLIENT_READY

ROOT = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Fidelity stage gate tests
# ---------------------------------------------------------------------------


class TestFidelityStageGate:
    def test_semantic_only_enabled_rejected(self):
        """semantic-only + enabled must fail the gate."""
        features = {"visual_fidelity": "semantic-only", "status": "enabled"}
        ok, msg = vf.check_fidelity_stage_gate(features)
        assert not ok
        assert "semantic-only" in msg

    def test_placeholder_enabled_rejected(self):
        """placeholder + enabled must fail the gate."""
        features = {"visual_fidelity": "placeholder", "status": "enabled"}
        ok, msg = vf.check_fidelity_stage_gate(features)
        assert not ok
        assert "placeholder" in msg

    def test_high_fidelity_enabled_accepted(self):
        """high-fidelity + enabled must pass the gate."""
        features = {"visual_fidelity": "high-fidelity", "status": "enabled"}
        ok, msg = vf.check_fidelity_stage_gate(features)
        assert ok

    def test_no_fidelity_field_skips(self):
        """No fidelity field → skip (pass)."""
        features = {"status": "enabled"}
        ok, msg = vf.check_fidelity_stage_gate(features)
        assert ok

    def test_semantic_only_planned_accepted(self):
        """semantic-only + planned is fine (not enabled)."""
        features = {"visual_fidelity": "semantic-only", "status": "planned"}
        ok, msg = vf.check_fidelity_stage_gate(features)
        assert ok


class TestRegistryFidelityGate:
    def test_empty_registry_passes(self):
        """Empty registry should have no violations."""
        violations = vf.check_registry_fidelity_gate({"entries": []})
        assert violations == []

    def test_good_registry_passes(self):
        """Registry without fidelity-enabled mismatch passes."""
        registry = {
            "entries": [
                {
                    "family": "test-family",
                    "graphical_variants": [
                        {
                            "graphical_variant": "good-variant",
                            "status": "enabled",
                            "features": {"visual_fidelity": "high-fidelity"},
                        },
                        {
                            "graphical_variant": "planned-variant",
                            "status": "planned",
                            "features": {"visual_fidelity": "placeholder"},
                        },
                    ],
                }
            ]
        }
        violations = vf.check_registry_fidelity_gate(registry)
        assert violations == []

    def test_bad_registry_detected(self):
        """Registry with semantic-only enabled is detected."""
        registry = {
            "entries": [
                {
                    "family": "test-family",
                    "graphical_variants": [
                        {
                            "graphical_variant": "bad-variant",
                            "status": "enabled",
                            "features": {"visual_fidelity": "semantic-only"},
                        }
                    ],
                }
            ]
        }
        violations = vf.check_registry_fidelity_gate(registry)
        assert len(violations) == 1
        assert "semantic-only" in violations[0]


# ---------------------------------------------------------------------------
# Verdict/report tests
# ---------------------------------------------------------------------------


class TestVerdict:
    def test_empty_verdict_ok(self):
        v = vf.VisualFidelityVerdict(0)
        assert v.all_passed
        assert v.failures == []

    def test_verdict_tracks_failures(self):
        v = vf.VisualFidelityVerdict(0)
        v.add("check1", True, "ok")
        v.add("check2", False, "failed")
        assert not v.all_passed
        assert len(v.failures) == 1
        assert v.failures[0].name == "check2"


class TestFidelityReport:
    def test_empty_report_ok(self):
        r = vf.FidelityReport()
        assert r.ok
        assert r.total_checks == 0
        assert r.total_failures == 0

    def test_report_aggregates(self):
        r = vf.FidelityReport()
        v1 = vf.VisualFidelityVerdict(0)
        v1.add("a", True)
        v2 = vf.VisualFidelityVerdict(1)
        v2.add("b", False)
        r.add_verdict(v1)
        r.add_verdict(v2)
        assert not r.ok
        assert r.total_checks == 2
        assert r.total_failures == 1


# ---------------------------------------------------------------------------
# PPTX-based measurable checks
# ---------------------------------------------------------------------------


class TestGraphicalAreaSufficient:
    def test_empty_slide_fails(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide = prs.slides[0]
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_graphical_area_sufficient(slide, 0, verdict, min_occupancy=0.05)
        assert not verdict.all_passed
        assert "graphical_area_sufficient" in [c.name for c in verdict.failures]
        assert any("0.000" in c.detail for c in verdict.failures)

    def test_large_shape_passes(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Add a large shape covering ~30% of slide area
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8), Inches(8)
        )
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_graphical_area_sufficient(slide, 0, verdict, min_occupancy=0.05)
        assert verdict.all_passed


class TestWhiteOnWhite:
    def test_white_fill_no_border_detected(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(3)
        )
        shape.name = "pattern:test:white-card"
        # Set white fill
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        # No border set — invisible
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_no_white_on_white(slide, 0, verdict)
        assert not verdict.all_passed
        assert "no_white_on_white" in [c.name for c in verdict.failures]

    def test_white_fill_with_border_passes(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(3)
        )
        shape.name = "pattern:test:white-card-with-border"
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Set border
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(1)
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_no_white_on_white(slide, 0, verdict)
        assert verdict.all_passed


class TestColorRoles:
    def test_no_colors_fails(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_color_roles_sufficient(slide, 0, verdict, required_roles=1)
        assert not verdict.all_passed

    def test_colors_present_passes(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 102, 204)
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_color_roles_sufficient(slide, 0, verdict, required_roles=1)
        assert verdict.all_passed


class TestTextToGraphicsRatio:
    def test_all_text_fails(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Add a text box (no fill, just text)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(8))
        tf = txBox.text_frame
        tf.text = "A" * 500
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_text_to_graphics_ratio(slide, 0, verdict, max_ratio=3.0)
        # With text only and no shapes, ratio will be very high
        assert not verdict.all_passed

    def test_mix_passes(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Add a large shape
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(6)
        )
        # Add a small text
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tf = txBox.text_frame
        tf.text = "Short text"
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_text_to_graphics_ratio(slide, 0, verdict, max_ratio=5.0)
        assert verdict.all_passed


class TestSpatialBalance:
    def test_single_shape_passes(self):
        """Single shape should skip (too few shapes)."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(2)
        )
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_spatial_balance(slide, 0, verdict)
        assert verdict.all_passed

    def test_balanced_shapes_passes(self):
        """Shapes on both sides should pass."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Two left, two right
        for x in [1, 3, 11, 13]:
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(1), Inches(2), Inches(2)
            )
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_spatial_balance(slide, 0, verdict)
        assert verdict.all_passed



class TestUnnaturalShortWrapping:
    """Tests for check_unnatural_short_wrapping.

    Operates on synthetic Presentation objects with text frames.
    """
    def test_short_last_line_detected(self):
        """Three paragraphs where the last is 'the' (3 chars) — must FAIL."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(5))
        tf = txBox.text_frame
        tf.clear()
        texts = [
            "This is a very long first paragraph that should be well above any threshold",
            "Another medium length second paragraph that continues the narrative flow",
            "the",  # very short last line
        ]
        for i, txt in enumerate(texts):
            if i == 0:
                tf.text = txt
            else:
                p = tf.add_paragraph()
                p.text = txt
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_unnatural_short_wrapping(slide, 0, verdict)
        # Detection branch: short last paragraph must be flagged
        assert not verdict.all_passed
        assert "unnatural_short_wrapping" in [c.name for c in verdict.failures]

    def test_three_long_paragraphs_passes(self):
        """Three paragraphs all of reasonable length — must PASS."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(5))
        tf = txBox.text_frame
        tf.clear()
        texts = [
            "First paragraph with enough content to be well above any short threshold",
            "Second paragraph that continues the narrative with similar length",
            "Third paragraph that is also long enough to not be flagged as short",
        ]
        for i, txt in enumerate(texts):
            if i == 0:
                tf.text = txt
            else:
                p = tf.add_paragraph()
                p.text = txt
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_unnatural_short_wrapping(slide, 0, verdict)
        assert verdict.all_passed


class TestRequiredIcons:
    """Tests for check_required_icons.

    Verifies icon counting for small filled shapes.
    """
    def test_small_filled_shapes_count_as_icons(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Add small filled shapes (icon-like)
        for i in range(3):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.2 + i * 1.5), Inches(0.5), Inches(0.6), Inches(0.6),
            )
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 102, 204)
            shape.name = f"icon:{i}"
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_required_icons(slide, 0, verdict, required_icon_count=3)
        assert verdict.all_passed

    def test_no_icons_when_required_passes(self):
        """required_icon_count=0 should skip the check."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_required_icons(slide, 0, verdict, required_icon_count=0)
        assert verdict.all_passed

    def test_icon_count_below_required_fails(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Only 1 icon but 3 required
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.7), Inches(0.7),
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 102, 204)
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_required_icons(slide, 0, verdict, required_icon_count=3)
        assert not verdict.all_passed


class TestRequiredLayers:
    """Tests for check_required_layers.

    Verifies visual layer detection heuristic.
    """
    def test_single_layer_when_required_one_passes(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # No shapes = effectively 1 layer
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_required_layers(slide, 0, verdict, required_layer_count=1)
        assert verdict.all_passed

    def test_multiple_layers_detected(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Large background shape (background layer)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(20), Inches(11),
        )
        bg_fill = bg.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(240, 240, 240)
        # Small card shape (card layer)
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2), Inches(2), Inches(5), Inches(4),
        )
        card_fill = card.fill
        card_fill.solid()
        card_fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Text box (text layer)
        txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(3), Inches(1))
        tf = txBox.text_frame
        tf.text = "Sample text"
        verdict = vf.VisualFidelityVerdict(0)
        vf.check_required_layers(slide, 0, verdict, required_layer_count=2)
        assert verdict.all_passed


class TestRealRegistryEnforcement:

    def test_all_enabled_variants_have_visual_fidelity(self):
        """Every enabled variant in the real registry must have visual_fidelity set.

        This test documents completeness — it should pass once all variants are classified.
        """
        registry_path = ROOT / "schemas" / "pattern-registry.yaml"
        from pathlib import Path
        registry_path = Path(__file__).resolve().parent.parent / "schemas" / "pattern-registry.yaml"
        with open(registry_path) as f:
            reg = yaml.safe_load(f)
        missing: list[str] = []
        for entry in reg.get("entries", []):
            family = entry.get("family", "?")
            for variant in entry.get("graphical_variants", []):
                if variant.get("status") == "enabled":
                    gv = variant.get("graphical_variant", "?")
                    features = variant.get("features", {})
                    if "visual_fidelity" not in features:
                        missing.append(f"{family}/{gv}")
        assert not missing, f"Enabled variants missing visual_fidelity: {missing}"

    def test_enabled_variants_not_placeholder(self):
        """Verify that NO enabled variant has visual_fidelity=placeholder in the real registry.

        Currently expected to FAIL because all 16 enabled variants are unclassified.
        This test serves as documentation of the human-required gap.
        Once a human classifies variants, this test should be adjusted or removed.
        """
        registry_path = ROOT / "schemas" / "pattern-registry.yaml"
        registry_path = Path(__file__).resolve().parent.parent / "schemas" / "pattern-registry.yaml"
        with open(registry_path) as f:
            reg = yaml.safe_load(f)
        violations: list[str] = []
        for entry in reg.get("entries", []):
            family = entry.get("family", "?")
            for variant in entry.get("graphical_variants", []):
                if variant.get("status") == "enabled":
                    gv = variant.get("graphical_variant", "?")
                    features = variant.get("features", {})
                    fidelity = features.get("visual_fidelity", "placeholder")
                    if fidelity in NON_CLIENT_READY:
                        violations.append(
                            f"{family}/{gv}: visual_fidelity={fidelity!r} but status=enabled"
                        )
        # This is expected to fail — documenting the real-world gap
        # Update once human classification is complete
        assert not violations, (
            f"All {len(violations)} enabled variants are placeholder/semantic-only.\n",
            "This is expected — human classification required. Violations:\n",
            "\n".join(violations),
        )
