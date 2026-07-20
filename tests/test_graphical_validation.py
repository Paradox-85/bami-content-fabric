"""Tests for the graphical/topology validator (shared/pptx/graphical_validation.py).

Verifies:
- Off-canvas pattern shape detection
- Funnel monotonic width detection
- Circle loop closure detection
- Step connector sequence detection
- Report behavior
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from shared.pptx import graphical_validation as gv


ROOT = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Report tests
# ---------------------------------------------------------------------------

class TestReport:
    def test_empty_report_is_ok(self):
        r = gv.Report()
        assert r.ok
        assert len(r.violations) == 0

    def test_report_adds_violation(self):
        r = gv.Report()
        r.add(0, "test violation")
        assert not r.ok
        assert len(r.violations) == 1
        assert "slide 0:" in r.violations[0]


# ---------------------------------------------------------------------------
# Off-canvas checks
# ---------------------------------------------------------------------------

class TestOffCanvas:
    def test_all_shapes_on_canvas(self):
        """Shapes within canvas bounds should not trigger violations."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide = prs.slides[0]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(3)
        )
        shape.name = "pattern:test:normal"
        rep = gv.Report()
        gv.check_no_off_canvas(slide, 0, rep)
        assert rep.ok

    def test_off_canvas_pattern_shape_detected(self):
        """Pattern shapes placed off-canvas should trigger a violation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        # Shape far to the left
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(-2), Inches(1), Inches(3), Inches(1)
        )
        shape.name = "pattern:test:off-left"
        rep = gv.Report()
        gv.check_no_off_canvas(slide, 0, rep)
        assert not rep.ok
        assert "off-canvas" in rep.violations[0]

    def test_off_canvas_non_pattern_ignored(self):
        """Non-pattern shapes off-canvas should NOT trigger a violation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(-2), Inches(1), Inches(3), Inches(1)
        )
        shape.name = "normal-shape"
        rep = gv.Report()
        gv.check_no_off_canvas(slide, 0, rep)
        assert rep.ok


# ---------------------------------------------------------------------------
# Funnel monotonic width
# ---------------------------------------------------------------------------

class TestFunnelMonotonic:
    def test_monotonic_narrowing_passes(self):
        """Funnel segments that monotonically narrow should pass."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        # Add shapes with decreasing widths (simulating a funnel)
        for i, w in enumerate([8.0, 6.0, 4.0, 2.0]):
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches((9 - w) / 2), Inches(0.5 + i * 1.0),
                Inches(w), Inches(0.6),
            )
            s.name = f"funnel_seg_{i}"

        rep = gv.Report()
        gv.check_funnel_monotonic_width(slide, 0, rep)
        assert rep.ok

    def test_non_monotonic_fails(self):
        """Non-monotonic funnel should trigger a violation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        # Add shapes with non-monotonic widths
        for i, w in enumerate([8.0, 6.0, 7.0, 2.0]):  # 7.0 > 6.0 = violation
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches((9 - w) / 2), Inches(0.5 + i * 1.0),
                Inches(w), Inches(0.6),
            )
            s.name = f"funnel_seg_{i}"

        rep = gv.Report()
        gv.check_funnel_monotonic_width(slide, 0, rep)
        assert not rep.ok
        assert "non-monotonic" in rep.violations[0].lower()


# ---------------------------------------------------------------------------
# Circle loop closure
# ---------------------------------------------------------------------------

class TestCircleLoopClosure:
    def test_sufficient_connectors_passes(self):
        """Circle-steps with enough connectors should pass."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        # Add circles
        for i in range(4):
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1 + i * 2), Inches(1), Inches(0.5), Inches(0.5)
            )
            c.name = f"pattern:circular-process-loop/circle-steps:node:{i:02d}:circle"

        # Add connectors (thin rectangles)
        for i in range(4):
            conn = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1 + i * 2 + 0.5), Inches(1.2), Inches(0.3), Inches(0.03)
            )
            conn.name = f"pattern:circular-process-loop/circle-steps:connector:{i:02d}"

        rep = gv.Report()
        gv.check_circle_loop_closure(slide, 0, rep)
        assert rep.ok

    def test_missing_connectors_fails(self):
        """Circle-steps without connectors should trigger a violation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        for i in range(4):
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1 + i * 2), Inches(1), Inches(0.5), Inches(0.5)
            )
            c.name = f"pattern:circular-process-loop/circle-steps:node:{i:02d}:circle"

        rep = gv.Report()
        gv.check_circle_loop_closure(slide, 0, rep)
        assert not rep.ok
        assert "connector" in rep.violations[0].lower()


# ---------------------------------------------------------------------------
# Step connector sequence
# ---------------------------------------------------------------------------

class TestStepConnectorSequence:
    def test_right_arrows_present(self):
        """Step sequence with right arrows should pass."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        for i in range(3):
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1 + i * 2), Inches(1), Inches(0.5), Inches(0.5)
            )
            c.name = f"step_circle_{i}"

        for i in range(2):
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(2 + i * 2), Inches(1.1), Inches(0.4), Inches(0.1)
            )
            a.name = f"step_arrow_{i}"

        rep = gv.Report()
        gv.check_step_connector_sequence(slide, 0, rep)
        assert rep.ok

    def test_wrong_arrow_count(self):
        """Wrong number of arrows should trigger a violation."""
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        for i in range(3):
            c = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(1 + i * 2), Inches(1), Inches(0.5), Inches(0.5)
            )
            c.name = f"step_circle_{i}"

        for i in range(1):  # Only 1 arrow for 3 circles — should be 2
            a = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(2), Inches(1.1), Inches(0.4), Inches(0.1)
            )
            a.name = f"step_arrow_{i}"

        rep = gv.Report()
        gv.check_step_connector_sequence(slide, 0, rep)
        assert not rep.ok
        assert "connector" in rep.violations[0].lower() or "arrow" in rep.violations[0].lower()


# ---------------------------------------------------------------------------
# Integration: validate on generated sample deck
# ---------------------------------------------------------------------------

class TestIntegration:
    """Integration tests that build a deck and validate it."""

    def test_validate_generated_deck(self, tmp_path, tmp_out, tokens_path, template_path):
        """Build a minimal deck and run graphical validation on the output."""
        from shared.pptx.build import build_deck

        deck = {
            "title": "Graphical Validation Test",
            "slides": [
                {"template": "cover", "fields": {"hero": "Test"}},
                {
                    "template": "content",
                    "fields": {"title": "Process"},
                    "content": {"items": ["Step A", "Step B", "Step C"]},
                },
                {"template": "closing", "fields": {}},
            ],
        }
        deck_path = tmp_path / "_graphical_test.json"
        deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

        result = build_deck(deck_path, tmp_out, template_path, tokens_path)
        assert result["slides_rendered"] == 3
        assert tmp_out.exists()

        rep = gv.validate(tmp_out)
        assert rep.ok, f"Graphical validation violations: {rep.violations}"
