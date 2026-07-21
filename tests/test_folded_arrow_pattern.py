"""Tests for the folded-arrow-horizontal native pattern injector.

Verifies:
- Injector produces correct shape count
- Shape naming convention is followed
- Content-only auto-resolution generates inject-pattern blocks
- Legacy explicit layout: numbered-process-steps still works
- Items+bodies legacy content is transformed
"""

from __future__ import annotations

import json

from pptx import Presentation

from shared.pptx.build import _legacy_content_to_steps
from shared.pptx.pattern_injectors.registry import get_injector, list_injectors
from shared.pptx.pattern_selection import resolve_pattern
from tools.pptx_validate.cli import validate

# ---------------------------------------------------------------------------
# Unit: injector registration
# ---------------------------------------------------------------------------

def test_folded_arrow_injector_registered():
    """The folded-arrow-horizontal injector is registered and callable."""
    assert "folded-arrow-horizontal" in list_injectors()
    fn = get_injector("folded-arrow-horizontal")
    assert fn is not None
    assert callable(fn)


# ---------------------------------------------------------------------------
# Unit: legacy content transformation
# ---------------------------------------------------------------------------

def test_legacy_items_to_steps():
    """Plain 'items' strings are transformed to steps with auto-numbering."""
    result = _legacy_content_to_steps({"items": ["A", "B", "C"]})
    assert len(result) == 3
    assert result[0] == {"number": "01", "title": "A"}
    assert result[1] == {"number": "02", "title": "B"}
    assert result[2] == {"number": "03", "title": "C"}


def test_legacy_items_with_bodies():
    """Parallel 'items' + 'bodies' arrays produce steps with body text."""
    result = _legacy_content_to_steps({
        "items": ["Plan", "Do"],
        "bodies": ["Define scope", "Execute"],
    })
    assert len(result) == 2
    assert result[0] == {"number": "01", "title": "Plan", "body": "Define scope"}
    assert result[1] == {"number": "02", "title": "Do", "body": "Execute"}


def test_legacy_steps_dicts():
    """'steps' as dicts with title/body are passed through."""
    result = _legacy_content_to_steps({
        "steps": [
            {"title": "Step 1", "body": "Body 1"},
            {"title": "Step 2"},
        ]
    })
    assert len(result) == 2
    assert result[0] == {"number": "01", "title": "Step 1", "body": "Body 1"}
    assert result[1] == {"number": "02", "title": "Step 2"}


def test_legacy_items_dicts():
    """'items' as dicts with title/body are passed through."""
    result = _legacy_content_to_steps({
        "items": [
            {"title": "Alpha", "body": "Desc A", "number": "1"},
            {"title": "Beta"},
        ]
    })
    assert len(result) == 2
    assert result[0] == {"number": "1", "title": "Alpha", "body": "Desc A"}
    assert result[1] == {"number": "02", "title": "Beta"}


# ---------------------------------------------------------------------------
# Unit: resolver enrichment
# ---------------------------------------------------------------------------

def test_resolver_returns_folded_arrow_variant():
    """Content with items → resolves to folded-arrow-horizontal variant."""
    class FakeTokens:
        def __init__(self, brand="bami"):
            self._brand = brand
        @property
        def raw(self):
            return {"brand": self._brand}

    content = {"items": ["A", "B", "C"]}
    tokens = FakeTokens("bami")
    sel = resolve_pattern(content, tokens)
    assert sel.family == "numbered-process-steps"
    assert sel.pattern_template_id == "numbered-process-steps/folded-arrow-horizontal@1.0.0"
    assert sel.graphical_variant == "folded-arrow-horizontal"
    assert sel.renderer_binding is not None
    assert sel.renderer_binding["native"]["injector_id"] == "folded-arrow-horizontal"
    assert sel.features is not None
    assert sel.features.get("native_editable") is True


# ---------------------------------------------------------------------------
# Integration: build pipeline with explicit injector block
# ---------------------------------------------------------------------------

def test_build_pipeline_folded_arrow_injector(tmp_path, tmp_out, tokens_path, template_path):
    """A deck with a content slide resolves through folded-arrow injector and builds."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Folded Arrow Build Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Process"},
                "content": {
                    "items": ["Step A", "Step B", "Step C"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_folded_arrow_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    # Validate the output PPTX
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: explicit layout still works
# ---------------------------------------------------------------------------

def test_explicit_layout_numbered_process_steps_still_works(tmp_path, tmp_out, tokens_path, template_path):
    """Explicit layout: numbered-process-steps still renders through legacy layout stub."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Explicit Layout Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Steps"},
                "layout": "numbered-process-steps",
                "content": {
                    "items": ["Alpha", "Beta", "Gamma"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_explicit_layout_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: explicit kind: steps block still works
# ---------------------------------------------------------------------------

def test_explicit_kind_steps_block_still_works(tmp_path, tmp_out, tokens_path, template_path):
    """Explicit kind: steps block still renders through add_steps builder."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Explicit Steps Block Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Steps Block"},
                "blocks": [
                    {
                        "kind": "steps",
                        "x": 0.6, "y": 1.5, "w": 18.8,
                        "count": 3,
                        "numbers": ["01", "02", "03"],
                        "titles": ["Step A", "Step B", "Step C"],
                    }
                ],
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_explicit_steps_block_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: shape naming convention in generated PPTX
# ---------------------------------------------------------------------------

def test_folded_arrow_shape_naming_convention(tmp_path, tmp_out, tokens_path, template_path):
    """The generated PPTX contains shapes with the deterministic naming convention."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "Shape Naming Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Named Process"},
                "content": {
                    "items": ["Alpha", "Beta", "Gamma"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_shape_naming_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    # Open the PPTX and inspect shape names on the content slide (index 1)
    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    shape_names = [shp.name for shp in slide.shapes]
    # Check naming convention pattern
    pattern_shapes = [n for n in shape_names if n.startswith("pattern:numbered-process-steps/folded-arrow-horizontal")]
    assert len(pattern_shapes) >= 3, (
        f"Expected at least 3 pattern-named shapes, got {len(pattern_shapes)}: {pattern_shapes}"
    )
    # Verify role naming (use endswith for exact role suffix, not substring)
    circle_shapes = [n for n in pattern_shapes if n.endswith(":circle")]
    number_shapes = [n for n in pattern_shapes if n.endswith(":number")]
    title_shapes = [n for n in pattern_shapes if n.endswith(":title")]
    connector_shapes = [n for n in pattern_shapes if ":connector:" in n]
    assert len(circle_shapes) == 3, f"Expected 3 circle shapes, got {len(circle_shapes)}"
    assert len(number_shapes) == 3, f"Expected 3 number shapes, got {len(number_shapes)}"
    assert len(title_shapes) == 3, f"Expected 3 title shapes, got {len(title_shapes)}"
    assert len(connector_shapes) == 2, f"Expected 2 connector shapes, got {len(connector_shapes)}"

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


# ---------------------------------------------------------------------------
# Integration: 5 items (max capacity) builds fine
# ---------------------------------------------------------------------------

def test_folded_arrow_5_items(tmp_path, tmp_out, tokens_path, template_path):
    """5 items resolve through folded-arrow and build successfully."""
    from shared.pptx.build import build_deck

    deck = {
        "title": "5 Steps Test",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Five Steps"},
                "content": {
                    "items": ["One", "Two", "Three", "Four", "Five"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = tmp_path / "_5_steps_test.json"
    deck_path.write_text(json.dumps(deck, indent=2), encoding="utf-8")

    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()

    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"
