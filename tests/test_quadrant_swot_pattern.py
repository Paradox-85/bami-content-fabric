"""Tests for the quadrant-swot (SWOT matrix) native PPTX injector."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from shared.pptx.tokens import load_tokens

ROOT = Path(__file__).resolve().parent.parent


def _load_tokens(tokens_path: Path):
    """Helper to load tokens from a path."""
    return load_tokens(tokens_path)


class TestQuadrantSwotInjector:
    def test_swot_injector_creates_4_labeled_quadrants(self, tmp_path, tokens_path):
        """SWOT injector should produce shapes for all 4 quadrants."""
        tokens = _load_tokens(tokens_path)

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("quadrant-swot")
        assert injector is not None, "quadrant-swot injector not registered"

        quadrants = [
            {"title": "Strong Brand", "body": "Market leader in 3 regions"},
            {"title": "High Costs", "body": "Operating margin below target"},
            {"title": "New Markets", "body": "APAC expansion underway"},
            {"title": "Competitors", "body": "3 new entrants this year"},
        ]

        shapes = injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0, quadrants=quadrants)
        assert len(shapes) >= 8, f"Expected >=8 shapes, got {len(shapes)}"

        # Check that shapes include at least 4 header bars
        header_count = 0
        for shp in shapes:
            if shp.has_text_frame:
                txt = shp.text_frame.text.strip()
                if txt in ("Strengths", "Weaknesses", "Opportunities", "Threats"):
                    header_count += 1
        assert header_count >= 4, f"Expected >=4 SWOT headers, got {header_count}"

    def test_swot_injector_raises_without_quadrants(self, tmp_path, tokens_path):
        """SWOT injector should raise ValueError when quadrants is missing."""
        tokens = _load_tokens(tokens_path)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("quadrant-swot")
        assert injector is not None

        with pytest.raises(ValueError, match="quadrant-swot: 'quadrants' parameter is required"):
            injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0)

    def test_swot_injector_handles_less_than_4_quadrants(self, tmp_path, tokens_path):
        """SWOT injector should pad to 4 quadrants when fewer are given."""
        tokens = _load_tokens(tokens_path)
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        slide = prs.slides[0]

        from shared.pptx.pattern_injectors.registry import get_injector

        injector = get_injector("quadrant-swot")
        assert injector is not None

        quadrants = [
            {"title": "Strength", "body": "Good"},
            {"title": "Weakness", "body": "Bad"},
        ]

        shapes = injector(slide, tokens, x=0.5, y=0.5, w=9.0, h=5.0, quadrants=quadrants)
        assert len(shapes) >= 8, f"Expected >=8 shapes, got {len(shapes)}"

    def test_swot_variant_registered_in_registry(self):
        """The quadrant-matrix/swot-grid variant should be registered in the pattern registry."""
        import yaml

        registry_path = ROOT / "schemas" / "pattern-registry.yaml"
        with registry_path.open(encoding="utf-8") as f:
            registry = yaml.safe_load(f)

        for entry in registry.get("entries", []):
            if entry.get("family") == "quadrant-matrix":
                variants = [v.get("graphical_variant") for v in entry.get("graphical_variants", [])]
                assert "swot-grid" in variants, "swot-grid variant not found in quadrant-matrix family"
                return
        pytest.fail("quadrant-matrix family not found in registry")
