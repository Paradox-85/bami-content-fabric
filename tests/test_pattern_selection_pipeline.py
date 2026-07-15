"""Integration tests: build pipeline with pattern selection resolver.

These tests verify that content-only slides (no explicit layout or blocks)
are resolved deterministically by the pattern resolver during ``build_deck``.
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from shared.pptx.build import build_deck
from shared.pptx.pattern_selection import resolve_pattern
from tools.pptx_validate.cli import validate
from tests.conftest import ROOT


FAKE_TOKENS_BAMI = ROOT / "templates" / "bami" / "design_tokens.yaml"


def _write_deck(tmp_path: Path, deck: dict) -> Path:
    path = tmp_path / "_content_deck.json"
    path.write_text(json.dumps(deck, indent=2), encoding="utf-8")
    return path


def test_content_kpis_resolves_via_pattern(tmp_path, tmp_out, tokens_path, template_path):
    """content with kpis but no layout → resolved to kpi_strip, build+validate ok."""
    deck = {
        "title": "Pattern-selected KPI deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Metrics"},
                "content": {
                    "kpis": [
                        {"number": "42", "label": "Units", "color": "primary"},
                        {"number": "99", "label": "Percent", "color": "positive"},
                    ],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify pattern selection produced expected layout
    content = deck["slides"][1]["content"]
    sel = resolve_pattern(content, tokens_path)
    assert sel.family == "kpi-dashboard-grid"
    assert sel.layout == "kpi_strip"
    # Validate the output PPTX
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


def test_content_periods_sections_resolves_via_pattern(tmp_path, tmp_out, tokens_path, template_path):
    """content with periods+sections but no layout → resolved to gantt, build+validate ok."""
    deck = {
        "title": "Pattern-selected Gantt deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Roadmap"},
                "content": {
                    "periods": [
                        {"label": "Q1", "key": "q1", "weeks": ["1", "2"]},
                        {"label": "Q2", "key": "q2", "weeks": ["3", "4"]},
                    ],
                    "tasks": [
                        {"label": "Init", "bars": [{"period_key": "q1", "start": 0.1, "duration": 0.8}]},
                    ],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify pattern selection
    content = deck["slides"][1]["content"]
    sel = resolve_pattern(content, tokens_path)
    assert sel.family == "gantt-matrix"
    assert sel.layout == "gantt"
    # Validate output
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


def test_content_unrecognizable_raises_build_error(tmp_path, tmp_out, tokens_path, template_path):
    """content with no recognizable keys → BuildError from resolver."""
    deck = {
        "title": "Unrecognizable content deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Weird"},
                "content": {"nonsense_key": "some_value"},
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    with pytest.raises(Exception):
        build_deck(deck_path, tmp_out, template_path, tokens_path)


def test_explicit_layout_still_wins(tmp_path, tmp_out, tokens_path, template_path):
    """explicit layout still overrides pattern resolver (backward compat)."""
    deck = {
        "title": "Explicit layout deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Compare"},
                "layout": "comparison_panel",
                "content": {
                    "kpis": [{"number": "42", "label": "Units"}],  # would be kpi_strip from resolver
                    "panels": [
                        {"title": "Option A", "heading": "A", "body": "Body A"},
                        {"title": "Option B", "heading": "B", "body": "Body B"},
                    ],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


def test_content_data_table_no_layout_materializes_block(tmp_path, tmp_out, tokens_path, template_path):
    """content-only slide with header+rows (data-table, layout:null) materializes a table block."""
    deck = {
        "title": "Data table deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Metrics Table"},
                "content": {
                    "header": ["Metric", "Value", "Status"],
                    "rows": [
                        ["Revenue", "$10M", "Green"],
                        ["Margin", "25%", "Yellow"],
                        ["Cost", "$5M", "Green"],
                    ],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify pattern selection
    sel = resolve_pattern(deck["slides"][1]["content"], tokens_path)
    assert sel.family == "data-table"
    assert sel.layout is None
    assert sel.block_kind == "table"
    # Validate the output PPTX builds successfully
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


def test_content_bullets_terminal_no_layout(tmp_path, tmp_out, tokens_path, template_path):
    """content-only slide with items=1 (below steps min) falls back to bullets, materializes block."""
    deck = {
        "title": "Single item deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Note"},
                "content": {
                    "items": ["Just one item"],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # items=1 fails min=2 for most families; should trigger fallback chain
    # We expect bullets (terminal fallback) to be selected
    sel = resolve_pattern(deck["slides"][1]["content"], tokens_path)
    assert sel.block_kind == "bullets"
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"


def test_content_impact_table_terminal_materializes_meaningfully(tmp_path, tmp_out, tokens_path, template_path):
    """content-only slide with rows (impact-table, layout:null) materializes a table with rows rendered."""
    from pptx import Presentation

    deck = {
        "title": "Impact table deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Risk Impact"},
                "content": {
                    "rows": [
                        ["Risk", "High"],
                        ["Cost", "Low"],
                        ["Timeline", "Medium"],
                    ],
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify pattern selection
    sel = resolve_pattern(deck["slides"][1]["content"], tokens_path)
    assert sel.family == "impact-table"
    assert sel.layout is None
    assert sel.block_kind == "table"
    # Validate the output PPTX: build must succeed without KeyError on header
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"
    # Inspect PPTX content: verify table cells contain input row values
    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    tbl_shapes = [s for s in slide.shapes if getattr(s, "has_table", False)]
    assert len(tbl_shapes) >= 1, "Expected at least one table shape on the content slide"
    tbl = tbl_shapes[0].table
    # Header row was auto-synthesized as Item | Value
    header_0 = tbl.cell(0, 0).text_frame.paragraphs[0].text.strip()
    assert header_0 == "Item", f"Expected header 'Item', got {header_0!r}"
    header_1 = tbl.cell(0, 1).text_frame.paragraphs[0].text.strip()
    assert header_1 == "Value", f"Expected header 'Value', got {header_1!r}"
    # Row 1: Risk | High
    assert tbl.cell(1, 0).text_frame.paragraphs[0].text.strip() == "Risk"
    assert tbl.cell(1, 1).text_frame.paragraphs[0].text.strip() == "High"
    # Row 2: Cost | Low
    assert tbl.cell(2, 0).text_frame.paragraphs[0].text.strip() == "Cost"
    assert tbl.cell(2, 1).text_frame.paragraphs[0].text.strip() == "Low"
    # Row 3: Timeline | Medium
    assert tbl.cell(3, 0).text_frame.paragraphs[0].text.strip() == "Timeline"
    assert tbl.cell(3, 1).text_frame.paragraphs[0].text.strip() == "Medium"


def test_content_before_after_split_terminal_materializes_meaningfully(tmp_path, tmp_out, tokens_path, template_path):
    """content-only slide with before+after (before-after-split, layout:null) materializes two darkcards."""
    from pptx import Presentation

    deck = {
        "title": "Before-after deck",
        "slides": [
            {"template": "cover", "fields": {"hero": "Test"}},
            {
                "template": "content",
                "fields": {"title": "Transformation"},
                "content": {
                    "before": "Legacy monolithic architecture with manual deployments and long release cycles.",
                    "after": "Microservices with CI/CD, automated testing, and weekly releases.",
                },
            },
            {"template": "closing", "fields": {}},
        ],
    }
    deck_path = _write_deck(tmp_path, deck)
    result = build_deck(deck_path, tmp_out, template_path, tokens_path)
    assert result["slides_rendered"] == 3
    assert tmp_out.exists()
    # Verify pattern selection
    sel = resolve_pattern(deck["slides"][1]["content"], tokens_path)
    assert sel.family == "before-after-split"
    assert sel.layout is None
    assert sel.block_kind == "darkcard"
    # Validate the output PPTX: build must succeed and not silently produce an empty slide
    rep = validate(tmp_out, tokens_path)
    assert rep.ok, f"Validation violations: {rep.violations}"
    # Inspect PPTX content: verify both before and after texts are present on the slide
    prs = Presentation(str(tmp_out))
    slide = prs.slides[1]
    all_text = []
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            for p in shp.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    all_text.append(txt)
    slide_text = " ".join(all_text)
    assert "Legacy monolithic" in slide_text, (
        f"Expected 'before' text in slide content, got: {slide_text[:200]}"
    )
    assert "Microservices with CI/CD" in slide_text, (
        f"Expected 'after' text in slide content, got: {slide_text[:200]}"
    )
